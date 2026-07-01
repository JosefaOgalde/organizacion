#!/usr/bin/env python3
"""Genera PPT: MOVA · GitHub Paso 1 — Crear cuenta."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "index/clientes/mkof/MOVA-GitHub-Paso1-Crear-Cuenta.pptx"

# Colores MOVA
C_BG = RGBColor(0xE8, 0xF6, 0xF8)
C_ACCENT = RGBColor(0x4A, 0x7A, 0x80)
C_ACCENT_DARK = RGBColor(0x2A, 0x4A, 0x4E)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT = RGBColor(0x1F, 0x23, 0x28)
C_MUTED = RGBColor(0x65, 0x6D, 0x76)
C_HIGHLIGHT = RGBColor(0xFF, 0xF8, 0xC5)
C_GITHUB_GREEN = RGBColor(0x1F, 0x88, 0x3D)
C_BLUE = RGBColor(0x09, 0x69, 0xDA)
C_BORDER = RGBColor(0xD0, 0xD7, 0xDE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_header_bar(slide, title, subtitle=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()
    add_textbox(slide, Inches(0.55), Inches(0.18), Inches(10), Inches(0.5), title, size=26, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.62), Inches(10), Inches(0.35), subtitle, size=13, color=C_BG)


def add_step_badge(slide, num, left=Inches(0.55), top=Inches(1.35)):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = C_ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER


def add_browser_mockup(slide, url, content_fn, left=Inches(6.8), top=Inches(1.5), width=Inches(5.8), height=Inches(5.2)):
    """Marco tipo navegador con barra de URL."""
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = C_WHITE
    frame.line.color.rgb = C_BORDER

    url_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.45))
    url_bar.fill.solid()
    url_bar.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFA)
    url_bar.line.color.rgb = C_BORDER
    add_textbox(slide, left + Inches(0.35), top + Inches(0.28), width - Inches(0.6), Inches(0.3), url, size=11, color=C_BLUE)

    body_left = left + Inches(0.35)
    body_top = top + Inches(0.85)
    body_width = width - Inches(0.7)
    content_fn(slide, body_left, body_top, body_width)


def mockup_field(slide, left, top, width, label, value, highlight=False):
    add_textbox(slide, left, top, width, Inches(0.25), label, size=11, bold=True, color=C_TEXT)
    field = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Inches(0.28), width, Inches(0.42))
    field.fill.solid()
    field.fill.fore_color.rgb = C_HIGHLIGHT if highlight else C_WHITE
    field.line.color.rgb = C_BLUE if highlight else C_BORDER
    add_textbox(slide, left + Inches(0.12), top + Inches(0.35), width - Inches(0.2), Inches(0.3), value, size=12, color=C_TEXT)


def mockup_button(slide, left, top, text, highlight=True):
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.2), Inches(0.48))
    btn.fill.solid()
    btn.fill.fore_color.rgb = C_GITHUB_GREEN
    btn.line.fill.background()
    tf = btn.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER
    if highlight:
        glow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left - Inches(0.05), top - Inches(0.05), Inches(2.3), Inches(0.58))
        glow.fill.background()
        glow.line.color.rgb = RGBColor(0x1A, 0x7F, 0x37)
        glow.line.width = Pt(2)
        # reorder - put glow behind (simplified: skip z-order, acceptable)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_ACCENT_DARK)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2), "MOVA · GitHub", size=44, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.8), "Paso 1: Crear cuenta con correo general", size=28, color=C_BG)
    add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.5), "Hito 1.1 · Respaldo n8n → GitHub  |  GRUPO MAKING OF", size=16, color=RGBColor(0xA8, 0xD8, 0xDC))
    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.4), "Guía para el encargado · Junio 2026", size=14, color=C_MUTED)


def slide_objetivo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "¿Qué vamos a hacer?", "Objetivo de esta guía")
    bullets = [
        "Crear una cuenta de GitHub con un correo general del proyecto (no personal).",
        "Esa cuenta servirá para respaldar los workflows de n8n (proyecto MOVA).",
        "Plan gratuito (Free) — suficiente para repos privados.",
        "Tiempo estimado: 10–15 minutos.",
        "Paso 2 (siguiente entrega): crear el repo privado mova-n8n-workflows.",
    ]
    y = Inches(1.5)
    for b in bullets:
        add_textbox(slide, Inches(0.9), y, Inches(11), Inches(0.55), f"•  {b}", size=20, color=C_TEXT)
        y += Inches(0.65)


def slide_correo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Antes de empezar: el correo", "Usar un buzón del equipo, no de una sola persona")

    ejemplos = ["infra@mova.cl", "tecnologia@grupomakingof.cl", "devops@mova.cl", "sistemas@empresa.cl"]
    x = Inches(0.7)
    for e in ejemplos:
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.45), Inches(2.6), Inches(0.45))
        tag.fill.solid()
        tag.fill.fore_color.rgb = C_WHITE
        tag.line.color.rgb = C_BORDER
        add_textbox(slide, x + Inches(0.1), Inches(1.52), Inches(2.4), Inches(0.35), e, size=12, color=C_ACCENT_DARK)
        x += Inches(2.75)

    reglas = [
        "Varias personas deben poder recibir el código de verificación.",
        "Guardar la contraseña en un gestor del equipo (1Password, Bitwarden…).",
        "Anotar usuario @ y correo en la ficha del proyecto MOVA.",
        "Activar verificación en dos pasos (2FA) cuando la cuenta esté lista.",
    ]
    y = Inches(2.3)
    for r in reglas:
        add_textbox(slide, Inches(0.9), y, Inches(11), Inches(0.5), f"✓  {r}", size=18, color=C_TEXT)
        y += Inches(0.55)

    warn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.8), Inches(11.5), Inches(1.6))
    warn.fill.solid()
    warn.fill.fore_color.rgb = C_HIGHLIGHT
    warn.line.color.rgb = RGBColor(0xD4, 0xA7, 0x2C)
    add_textbox(slide, Inches(0.95), Inches(5.0), Inches(11), Inches(1.2),
                "⚠  Si ya tienes otra sesión de GitHub abierta, usa ventana de incógnito\n     o cierra sesión antes de crear esta cuenta.", size=17, color=C_ACCENT_DARK)


def slide_paso(prs, num, titulo, texto, tip, url, mockup_type):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, f"Paso {num} de 9", titulo)
    add_step_badge(slide, num)

    add_textbox(slide, Inches(1.25), Inches(1.25), Inches(5.2), Inches(0.55), titulo, size=22, bold=True, color=C_ACCENT_DARK)
    add_textbox(slide, Inches(0.75), Inches(1.95), Inches(5.5), Inches(1.8), texto, size=17, color=C_TEXT)

    tip_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.0), Inches(5.6), Inches(1.5))
    tip_box.fill.solid()
    tip_box.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFA)
    tip_box.line.color.rgb = C_BORDER
    add_textbox(slide, Inches(0.9), Inches(4.15), Inches(5.2), Inches(1.2), f"💡 Tip: {tip}", size=14, color=C_MUTED)

    def content_signup_url(s, l, t, w):
        add_textbox(s, l, t, w, Inches(0.4), "GitHub · Sign up", size=16, bold=True)
        add_textbox(s, l, t + Inches(0.5), w, Inches(0.35), "Crear tu cuenta", size=12, color=C_MUTED)
        hint = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t + Inches(1.2), w, Inches(0.7))
        hint.fill.solid()
        hint.fill.fore_color.rgb = C_HIGHLIGHT
        hint.line.color.rgb = RGBColor(0xD4, 0xA7, 0x2C)
        add_textbox(s, l + Inches(0.15), t + Inches(1.35), w - Inches(0.3), Inches(0.5),
                    "👆 Escribe github.com/signup en la barra del navegador", size=12, color=C_ACCENT_DARK)

    def content_email(s, l, t, w):
        add_textbox(s, l, t, w, Inches(0.35), "GitHub", size=16, bold=True)
        mockup_field(s, l, t + Inches(0.5), w, "Email  ← AQUÍ", "infra@mova.cl", highlight=True)
        mockup_field(s, l, t + Inches(1.35), w, "Password", "••••••••••••")

    def content_password(s, l, t, w):
        mockup_field(s, l, t, w, "Email", "infra@mova.cl")
        mockup_field(s, l, t + Inches(0.85), w, "Password  ← AQUÍ", "Mova-2026-Segura!", highlight=True)
        add_textbox(s, l, t + Inches(1.75), w, Inches(0.3), "●●●●● Fuerte", size=11, color=C_GITHUB_GREEN)

    def content_username(s, l, t, w):
        mockup_field(s, l, t, w, "Username  ← AQUÍ", "mova-infra", highlight=True)
        add_textbox(s, l, t + Inches(0.95), w, Inches(0.3), "github.com/mova-infra", size=11, color=C_MUTED)

    def content_submit(s, l, t, w):
        cap = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(0.65))
        cap.fill.solid()
        cap.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFA)
        cap.line.color.rgb = C_BORDER
        add_textbox(s, l + Inches(0.15), t + Inches(0.15), w, Inches(0.4), "🤖 Verificar que no eres un robot", size=12, align=PP_ALIGN.CENTER)
        mockup_button(s, l + Inches(0.8), t + Inches(0.9), "Create account")

    def content_verify(s, l, t, w):
        add_textbox(s, l, t, w, Inches(0.35), "Verify your email", size=14, bold=True)
        add_textbox(s, l, t + Inches(0.4), w, Inches(0.35), "Código enviado a infra@mova.cl", size=11, color=C_MUTED)
        codes = "4   8   2   9   1   5   7   3"
        code_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t + Inches(0.9), w, Inches(0.55))
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = C_HIGHLIGHT
        code_box.line.color.rgb = C_BLUE
        add_textbox(s, l + Inches(0.2), t + Inches(1.0), w, Inches(0.4), codes, size=16, bold=True, align=PP_ALIGN.CENTER)
        mockup_button(s, l + Inches(0.9), t + Inches(1.65), "Continue", highlight=False)

    def content_skip(s, l, t, w):
        add_textbox(s, l, t, w, Inches(0.35), "Welcome to GitHub!", size=14, bold=True)
        add_textbox(s, l, t + Inches(0.45), w, Inches(0.35), "¿Cuántas personas trabajan contigo?", size=11, color=C_MUTED)
        btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l + Inches(0.5), t + Inches(1.1), Inches(2.5), Inches(0.45))
        btn.fill.background()
        btn.line.color.rgb = C_BORDER
        add_textbox(s, l + Inches(0.65), t + Inches(1.18), Inches(2.2), Inches(0.35), "Skip this step →", size=12, color=C_MUTED)

    def content_plan(s, l, t, w):
        add_textbox(s, l, t, w, Inches(0.35), "Pick a plan", size=14, bold=True)
        plan = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t + Inches(0.5), w, Inches(0.9))
        plan.fill.solid()
        plan.fill.fore_color.rgb = RGBColor(0xF0, 0xF6, 0xFF)
        plan.line.color.rgb = C_BLUE
        add_textbox(s, l + Inches(0.15), t + Inches(0.6), w, Inches(0.3), "Free", size=13, bold=True)
        add_textbox(s, l + Inches(0.15), t + Inches(0.9), w, Inches(0.3), "Repos privados ilimitados", size=11, color=C_MUTED)
        mockup_button(s, l + Inches(0.6), t + Inches(1.6), "Continue for free")

    def content_ready(s, l, t, w):
        add_textbox(s, l, t, w, Inches(0.5), "✓", size=36, bold=True, color=C_GITHUB_GREEN, align=PP_ALIGN.CENTER)
        add_textbox(s, l, t + Inches(0.55), w, Inches(0.35), "¡Cuenta creada!", size=14, bold=True, color=C_GITHUB_GREEN, align=PP_ALIGN.CENTER)
        dash = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t + Inches(1.1), w, Inches(0.75))
        dash.fill.solid()
        dash.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFA)
        dash.line.color.rgb = C_BORDER
        avatar = s.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.2), t + Inches(1.25), Inches(0.45), Inches(0.45))
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = C_ACCENT
        avatar.line.fill.background()
        add_textbox(s, l + Inches(0.8), t + Inches(1.2), w, Inches(0.3), "mova-infra", size=12, bold=True)
        add_textbox(s, l + Inches(0.8), t + Inches(1.5), w, Inches(0.3), "infra@mova.cl", size=10, color=C_MUTED)

    mockups = {
        "signup-url": content_signup_url,
        "signup-email": content_email,
        "signup-password": content_password,
        "signup-username": content_username,
        "signup-submit": content_submit,
        "verify-email": content_verify,
        "onboarding-skip": content_skip,
        "plan-free": content_plan,
        "dashboard-ready": content_ready,
    }

    mockup_fn = mockups.get(mockup_type, content_signup_url)
    add_browser_mockup(slide, url, mockup_fn)


def slide_checklist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Checklist final", "Marca cada ítem antes de avisar que terminaste")

    items = [
        "Correo general del proyecto definido y accesible por el equipo",
        "Cuenta creada en github.com/signup",
        "Correo verificado con el código de 8 dígitos",
        "Plan Free activo",
        "Usuario y contraseña guardados en gestor del equipo",
        "Datos anotados en ficha MOVA (correo + usuario @)",
    ]
    y = Inches(1.5)
    for item in items:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(0.35), Inches(0.35))
        box.fill.background()
        box.line.color.rgb = C_ACCENT
        box.line.width = Pt(2)
        add_textbox(slide, Inches(1.35), y - Inches(0.05), Inches(10.5), Inches(0.45), item, size=18, color=C_TEXT)
        y += Inches(0.65)


def slide_paso2_pendiente(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_ACCENT_DARK)
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(0.7), "Paso 2 · Próxima entrega", size=32, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.2),
                "Crear repositorio privado:\n\nmova-n8n-workflows", size=28, color=C_BG)
    add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.8),
                "Cuando la cuenta esté lista, avisar al equipo para recibir la guía del Paso 2.", size=18, color=RGBColor(0xA8, 0xD8, 0xDC))
    add_textbox(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4), "MOVA · GRUPO MAKING OF · Jun 2026", size=14, color=C_MUTED)


PASOS = [
    (1, "Abrir la página de registro", "En el navegador (Chrome, Edge o Firefox), escribe la dirección y presiona Enter.", "Si ya tienes sesión en GitHub con otra cuenta, abre una ventana de incógnito o cierra sesión primero.", "github.com/signup", "signup-url"),
    (2, "Ingresar el correo general", "En el campo Email, escribe el correo del proyecto (no un correo personal de una sola persona).", "Ejemplo: infra@mova.cl — debe ser un buzón al que el equipo tenga acceso.", "github.com/signup", "signup-email"),
    (3, "Crear una contraseña segura", "Elige una contraseña larga (mínimo 15 caracteres). Guárdala en el gestor de contraseñas del equipo.", "No uses la misma contraseña de otro servicio. GitHub mostrará si es débil o fuerte.", "github.com/signup", "signup-password"),
    (4, "Elegir nombre de usuario", "Este será el @público de la cuenta. Usa algo del proyecto, por ejemplo mova-infra o makingof-mova.", "Solo letras, números y guiones. Una vez creado, cuesta cambiarlo.", "github.com/signup", "signup-username"),
    (5, "Confirmar y crear cuenta", "GitHub puede pedir un puzzle o captcha. Complétalo y haz clic en Create account.", "Si el correo ya está usado, recupera acceso o usa otro correo general.", "github.com/signup", "signup-submit"),
    (6, "Verificar el correo", "GitHub envía un código de 8 dígitos al correo. Ábrelo, copia el código y pégalo en la pantalla.", "Revisa carpeta Spam si no llega en 2 minutos.", "github.com/signup", "verify-email"),
    (7, "Omitir preguntas opcionales", "GitHub puede preguntar cuántos integrantes tiene el equipo. Puedes saltar (Skip) todo.", "No afecta el uso. Lo importante es llegar al dashboard principal.", "github.com", "onboarding-skip"),
    (8, "Elegir plan gratuito", "Selecciona Continue for free. No necesitas plan de pago para repos privados.", "GitHub Free permite repos privados ilimitados — suficiente para MOVA.", "github.com", "plan-free"),
    (9, "Cuenta lista — anotar datos", "Ya estás dentro. Anota en la ficha MOVA: correo, usuario @, y dónde guardaste la contraseña.", "Siguiente paso: crear el repo privado mova-n8n-workflows.", "github.com", "dashboard-ready"),
]


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_objetivo(prs)
    slide_correo(prs)
    for p in PASOS:
        slide_paso(prs, *p)
    slide_checklist(prs)
    slide_paso2_pendiente(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Generado: {OUT}")
    print(f"Diapositivas: {len(prs.slides)}")


if __name__ == "__main__":
    main()
