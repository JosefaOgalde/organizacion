#!/usr/bin/env python3
"""Genera PPT: MOVA · GitHub Paso 2 — Repo privado mova-n8n-workflows."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "index/clientes/mkof/MOVA-GitHub-Paso2-Repo-Privado.pptx"

C_BG = RGBColor(0xE8, 0xF6, 0xF8)
C_ACCENT = RGBColor(0x4A, 0x7A, 0x80)
C_ACCENT_DARK = RGBColor(0x2A, 0x4A, 0x4E)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT = RGBColor(0x1F, 0x23, 0x28)
C_MUTED = RGBColor(0x65, 0x6D, 0x76)
C_HIGHLIGHT = RGBColor(0xFF, 0xF8, 0xC5)
C_GITHUB_GREEN = RGBColor(0x1F, 0x88, 0x3D)
C_BLUE = RGBColor(0x09, 0x69, 0xDA)
C_BORDER = RGBColor(0xD0, 0xD7, 0xDE)
C_PRIVATE = RGBColor(0x6E, 0x40, 0xC9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
REPO = "mova-n8n-workflows"


def set_slide_bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_header_bar(slide, title, subtitle=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()
    add_textbox(slide, Inches(0.55), Inches(0.18), Inches(10), Inches(0.5), title, size=26, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.62), Inches(10), Inches(0.35), subtitle, size=13, color=C_BG)


def add_step_badge(slide, num, left=Inches(0.55), top=Inches(1.35)):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = C_ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER


def mockup_field(slide, left, top, width, label, value, highlight=False):
    add_textbox(slide, left, top, width, Inches(0.25), label, size=11, bold=True, color=C_TEXT)
    field = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + Inches(0.28), width, Inches(0.42))
    field.fill.solid()
    field.fill.fore_color.rgb = C_HIGHLIGHT if highlight else C_WHITE
    field.line.color.rgb = C_BLUE if highlight else C_BORDER
    add_textbox(slide, left + Inches(0.12), top + Inches(0.35), width - Inches(0.2), Inches(0.3), value, size=12, color=C_TEXT)


def mockup_button(slide, left, top, text):
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.6), Inches(0.48))
    btn.fill.solid()
    btn.fill.fore_color.rgb = C_GITHUB_GREEN
    btn.line.fill.background()
    tf = btn.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER


def add_browser_mockup(slide, url, content_fn):
    left, top, width, height = Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2)
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = C_WHITE
    frame.line.color.rgb = C_BORDER
    url_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.45))
    url_bar.fill.solid()
    url_bar.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFA)
    url_bar.line.color.rgb = C_BORDER
    add_textbox(slide, left + Inches(0.35), top + Inches(0.28), width - Inches(0.6), Inches(0.3), url, size=11, color=C_BLUE)
    content_fn(slide, left + Inches(0.35), top + Inches(0.85), width - Inches(0.7))


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_ACCENT_DARK)
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.6), "Continuación · Paso 1 completado", size=18, color=RGBColor(0xA8, 0xD8, 0xDC))
    add_textbox(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(1.0), "MOVA · GitHub Paso 2", size=44, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.8), f"Crear repositorio privado: {REPO}", size=28, color=C_BG)
    add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5), "Hito 1.1 · Respaldo n8n → GitHub  |  GRUPO MAKING OF", size=16, color=RGBColor(0xA8, 0xD8, 0xDC))
    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.4), "Guía para el encargado · Junio 2026", size=14, color=C_MUTED)


def slide_objetivo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "¿Qué vamos a hacer?", "Paso 2 — después de crear la cuenta")
    bullets = [
        "Crear el repo privado mova-n8n-workflows en la cuenta del Paso 1.",
        "Visibilidad Private — los workflows pueden tener datos sensibles.",
        "Repo vacío (sin README) para el primer backup desde n8n.",
        "Anotar la URL del repo en la ficha MOVA.",
        "Tiempo estimado: 5–10 minutos.",
    ]
    y = Inches(1.5)
    for b in bullets:
        add_textbox(slide, Inches(0.9), y, Inches(11), Inches(0.55), f"•  {b}", size=20, color=C_TEXT)
        y += Inches(0.65)


def slide_requisitos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Requisitos", "Paso 1 debe estar listo")
    items = [
        "Cuenta GitHub con correo general del equipo.",
        "Usuario y contraseña en el gestor del equipo.",
        "Iniciar sesión en github.com/login.",
        f"Nombre del repo: {REPO} (exacto, minúsculas).",
    ]
    y = Inches(1.5)
    for item in items:
        add_textbox(slide, Inches(0.9), y, Inches(11), Inches(0.5), f"✓  {item}", size=18, color=C_TEXT)
        y += Inches(0.55)


def slide_paso(prs, num, titulo, texto, tip, url, mockup_type):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, f"Paso {num} de 10", titulo)
    add_step_badge(slide, num)
    add_textbox(slide, Inches(1.25), Inches(1.25), Inches(5.2), Inches(0.55), titulo, size=22, bold=True, color=C_ACCENT_DARK)
    add_textbox(slide, Inches(0.75), Inches(1.95), Inches(5.5), Inches(1.8), texto, size=17, color=C_TEXT)
    tip_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.0), Inches(5.6), Inches(1.5))
    tip_box.fill.solid()
    tip_box.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFA)
    tip_box.line.color.rgb = C_BORDER
    add_textbox(slide, Inches(0.9), Inches(4.15), Inches(5.2), Inches(1.2), f"💡 Tip: {tip}", size=14, color=C_MUTED)

    def login(s, l, t, w):
        mockup_field(s, l, t, w, "Username or email", "infra@mova.cl")
        mockup_field(s, l, t + Inches(0.85), w, "Password", "••••••••", highlight=True)
        mockup_button(s, l + Inches(0.5), t + Inches(1.75), "Sign in")

    def new_menu(s, l, t, w):
        add_textbox(s, l, t, w, Inches(0.35), "+  New repository", size=13, bold=True, color=C_BLUE)
        add_textbox(s, l, t + Inches(0.5), w, Inches(0.35), "github.com/new", size=11, color=C_MUTED)

    def owner(s, l, t, w):
        mockup_field(s, l, t, w, "Owner", "mova-infra", highlight=True)

    def repo_name(s, l, t, w):
        mockup_field(s, l, t, w, "Repository name", REPO, highlight=True)
        add_textbox(s, l, t + Inches(0.95), w, Inches(0.3), "✓ available", size=11, color=C_GITHUB_GREEN)

    def description(s, l, t, w):
        mockup_field(s, l, t, w, "Description", "Respaldo workflows n8n MOVA", highlight=True)

    def private(s, l, t, w):
        pub = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, Inches(0.45))
        pub.fill.solid()
        pub.fill.fore_color.rgb = C_WHITE
        pub.line.color.rgb = C_BORDER
        add_textbox(s, l + Inches(0.15), t + Inches(0.08), w, Inches(0.3), "○ Public", size=11)
        prv = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t + Inches(0.55), w, Inches(0.55))
        prv.fill.solid()
        prv.fill.fore_color.rgb = RGBColor(0xF6, 0xF0, 0xFF)
        prv.line.color.rgb = C_PRIVATE
        add_textbox(s, l + Inches(0.15), t + Inches(0.65), w, Inches(0.35), "● Private  🔒", size=12, bold=True, color=C_PRIVATE)

    def no_readme(s, l, t, w):
        add_textbox(s, l, t, w, Inches(0.35), "Initialize this repository with:", size=11, bold=True)
        for i, label in enumerate(["Add a README file", "Add .gitignore", "Choose a license"]):
            add_textbox(s, l + Inches(0.2), t + Inches(0.45 + i * 0.35), w, Inches(0.3), f"☐  {label}", size=10, color=C_MUTED)

    def create_btn(s, l, t, w):
        mockup_button(s, l + Inches(0.3), t + Inches(0.5), "Create repository")

    def repo_url(s, l, t, w):
        add_textbox(s, l, t, w, Inches(0.35), REPO, size=14, bold=True)
        add_textbox(s, l, t + Inches(0.45), w, Inches(0.35), "Private", size=11, color=C_PRIVATE)
        code = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t + Inches(1.0), w, Inches(0.55))
        code.fill.solid()
        code.fill.fore_color.rgb = C_HIGHLIGHT
        code.line.color.rgb = C_BLUE
        add_textbox(s, l + Inches(0.1), t + Inches(1.1), w, Inches(0.4), f"github.com/mova-infra/{REPO}", size=10)

    def collaborators(s, l, t, w):
        add_textbox(s, l, t, w, Inches(0.35), "Settings → Collaborators", size=12, bold=True)
        add_textbox(s, l, t + Inches(0.45), w, Inches(0.35), "Add people → correo del equipo", size=11, color=C_MUTED)

    mockups = {
        "login": login, "new-menu": new_menu, "owner": owner, "repo-name": repo_name,
        "description": description, "private": private, "no-readme": no_readme,
        "create-btn": create_btn, "repo-url": repo_url, "collaborators": collaborators,
    }
    add_browser_mockup(slide, url, mockups.get(mockup_type, new_menu))


def slide_checklist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Checklist final", "Antes de avisar al equipo técnico")
    items = [
        f"Repo {REPO} creado",
        "Visibilidad Private confirmada",
        "Sin README inicial",
        "URL anotada en ficha MOVA",
        "Colaboradores invitados (si aplica)",
    ]
    y = Inches(1.5)
    for item in items:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(0.35), Inches(0.35))
        box.fill.background()
        box.line.color.rgb = C_ACCENT
        box.line.width = Pt(2)
        add_textbox(slide, Inches(1.35), y - Inches(0.05), Inches(10.5), Inches(0.45), item, size=18, color=C_TEXT)
        y += Inches(0.65)


PASOS = [
    (1, "Iniciar sesión", "Entra en github.com/login con la cuenta del Paso 1.", "Usa el correo general del equipo.", "github.com/login", "login"),
    (2, "Abrir New repository", "Clic en + arriba a la derecha → New repository.", "También: github.com/new", "github.com/new", "new-menu"),
    (3, "Propietario (Owner)", "Verifica que Owner sea la cuenta MOVA.", "No cambies a otra org sin consultar.", "github.com/new", "owner"),
    (4, "Nombre del repo", f"Escribe exactamente: {REPO}", "Minúsculas y guiones.", "github.com/new", "repo-name"),
    (5, "Descripción", "Texto opcional para identificar el proyecto.", "Ej: Respaldo workflows n8n MOVA.", "github.com/new", "description"),
    (6, "Marcar Private", "Selecciona Private — nunca Public.", "Los workflows pueden tener secretos.", "github.com/new", "private"),
    (7, "Sin README inicial", "No marques README, .gitignore ni license.", "El primer push vendrá de n8n.", "github.com/new", "no-readme"),
    (8, "Create repository", "Revisa todo y pulsa el botón verde.", "Nombre + Private + vacío.", "github.com/new", "create-btn"),
    (9, "Copiar URL", "Anota la URL HTTPS del repo en ficha MOVA.", "La usará el backup de n8n.", f"github.com/.../{REPO}", "repo-url"),
    (10, "Colaboradores (opc.)", "Settings → Collaborators → Add people.", "Solo correos del equipo.", "github.com", "collaborators"),
]


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_title(prs)
    slide_objetivo(prs)
    slide_requisitos(prs)
    for p in PASOS:
        slide_paso(prs, *p)
    slide_checklist(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Generado: {OUT}")
    print(f"Diapositivas: {len(prs.slides)}")


if __name__ == "__main__":
    main()
