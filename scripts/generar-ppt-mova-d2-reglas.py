#!/usr/bin/env python3
"""Genera PPT: MOVA · Día 2 — Reglas mova_auth (auditoría + mapas conceptuales)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "index/clientes/mkof/MOVA-D2-Reglas-mova_auth.pptx"

C_BG = RGBColor(0xE8, 0xF6, 0xF8)
C_ACCENT = RGBColor(0x4A, 0x7A, 0x80)
C_ACCENT_DARK = RGBColor(0x2A, 0x4A, 0x4E)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT = RGBColor(0x1F, 0x23, 0x28)
C_MUTED = RGBColor(0x65, 0x6D, 0x76)
C_OK = RGBColor(0x1A, 0x7F, 0x37)
C_WARN = RGBColor(0xBF, 0x3F, 0x00)
C_HIGHLIGHT = RGBColor(0xFF, 0xF8, 0xC5)
C_PURPLE = RGBColor(0x6F, 0x42, 0xC1)
C_ORANGE_BG = RGBColor(0xFF, 0xE8, 0xD6)
C_GREEN_BG = RGBColor(0xD4, 0xED, 0xDA)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MODULOS = [
    ("Portal MOVA", "/mova/", "Google OAuth", "Sí", "Hoy no pasa"),
    ("mova_auth", "/mova_auth/", "PHP correo+clave", "Es el gate", "Parcial"),
    ("MOVA ERP", "/mova/erp/", "Contraseña local", "Sí", "Login aparte"),
    ("AXON", "/axon/", "Sin login visible", "Sí", "Validar D5"),
    ("RRHH", "/rrhh/", "403 Forbidden", "Sí", "Sin acceso público"),
    ("Documentos", "/documentos/", "Público", "No", "Excepción"),
]

CHECKLIST = [
    ("Regla de oro acordada con equipo técnico", False),
    ("Tabla de módulos revisada", False),
    ("Excepciones públicas confirmadas", False),
    ("Documento compartido (PDF/PPT)", False),
    ("Tarea mkof/02 marcada completada", False),
]


def set_slide_bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_header_bar(slide, title, subtitle="", accent=C_ACCENT):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_textbox(slide, Inches(0.55), Inches(0.18), Inches(11), Inches(0.5), title, size=24, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.62), Inches(11), Inches(0.35), subtitle, size=13, color=C_BG)


def slide_puente(prs, de_titulo, a_titulo, cuerpo, mapa_lineas=None):
    """Slide de transición entre bloques: texto explicativo + mini mapa ASCII."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, RGBColor(0xF3, 0xEE, 0xF8))
    add_header_bar(slide, "Puente — para entender el hilo", f"De «{de_titulo}» a «{a_titulo}»", C_PURPLE)

    badge = add_round_box(slide, Inches(0.75), Inches(1.25), Inches(1.35), Inches(0.42), C_PURPLE)
    add_textbox(slide, Inches(0.75), Inches(1.32), Inches(1.35), Inches(0.3), "PUENTE", size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2.25), Inches(1.28), Inches(10.3), Inches(0.45), cuerpo, size=15, color=C_TEXT)

    if mapa_lineas:
        box = add_round_box(slide, Inches(0.75), Inches(2.05), Inches(11.8), Inches(2.35), C_WHITE, C_PURPLE, 2)
        add_textbox(slide, Inches(0.95), Inches(2.2), Inches(2.5), Inches(0.3), "Mapa conceptual", size=12, bold=True, color=C_PURPLE)
        y = Inches(2.55)
        for line in mapa_lineas:
            add_textbox(slide, Inches(1.0), y, Inches(11.3), Inches(0.38), line, size=13, color=C_TEXT)
            y += Inches(0.42)

    flecha = add_round_box(slide, Inches(0.75), Inches(4.65), Inches(11.8), Inches(0.75), C_HIGHLIGHT, C_ACCENT)
    add_textbox(slide, Inches(0.95), Inches(4.82), Inches(11.4), Inches(0.45),
                f"Siguiente slide: {a_titulo}", size=14, bold=True, color=C_ACCENT_DARK, align=PP_ALIGN.CENTER)


def add_round_box(slide, left, top, width, height, fill, line=None, lw=1):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line:
        box.line.color.rgb = line
        box.line.width = Pt(lw)
    else:
        box.line.fill.background()
    return box


def add_arrow(slide, left, top, width=Inches(0.45), height=Inches(0.3), color=C_ACCENT):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    return arr


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_ACCENT_DARK)
    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5), "MOVA · Día 2 — Auditoría", size=20, color=RGBColor(0xA8, 0xD8, 0xDC))
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0), "Reglas mova_auth", size=44, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.7), "acme-chile.cl · único validador de sesión", size=22, color=C_BG)
    add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.5), "Tarea organizador: index.html?tarea=mkof/02", size=16, color=RGBColor(0xA8, 0xD8, 0xDC))
    add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.5), "8 jul 2026 · GRUPO MAKING OF", size=14, color=C_MUTED)
    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.6), "Solo documentar y acordar — NO tocar código en servidor", size=17, bold=True, color=C_HIGHLIGHT)


def slide_contexto(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Contexto — hallazgo Día 1", "Inventario completado · login fragmentado")
    items = [
        "Portal /mova/ → Google OAuth (no pasa por mova_auth).",
        "mova_auth/login.php → correo+clave (parcial, no gate del panel).",
        "mova/erp/ → contraseña local independiente.",
        "Ningún módulo M usa mova_auth como único validador.",
        "Antes de implementar (D3+), el equipo debe acordar las reglas por escrito.",
    ]
    y = Inches(1.5)
    for b in items:
        add_textbox(slide, Inches(0.85), y, Inches(11.5), Inches(0.55), f"•  {b}", size=17, color=C_TEXT)
        y += Inches(0.58)


def slide_mapa_fragmentado(prs):
    """Mapa conceptual: hoy hay varias puertas de entrada."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Mapa conceptual — situación HOY", "¿Por qué necesitamos reglas?")

    add_textbox(slide, Inches(0.75), Inches(1.25), Inches(12), Inches(0.45),
                "El usuario no sabe por dónde «entrar». Cada módulo decidió su propio login.", size=15, color=C_TEXT)

    user = add_round_box(slide, Inches(5.9), Inches(1.75), Inches(1.5), Inches(0.5), C_ACCENT_DARK)
    add_textbox(slide, Inches(5.9), Inches(1.85), Inches(1.5), Inches(0.35), "Usuario", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    puertas = [
        ("Google OAuth", "/mova/", "Entra con cuenta Google\nal panel MOVA", 0.7, C_ORANGE_BG, C_WARN),
        ("correo + clave", "mova_auth/login.php", "Login PHP aparte\nNO abre el panel", 3.35, C_HIGHLIGHT, C_ACCENT),
        ("contraseña local", "/mova/erp/", "Solo campo clave\nen el ERP", 6.0, C_ORANGE_BG, C_WARN),
        ("sin login visible", "/axon/", "Abre directo\n¿quién valida?", 8.65, RGBColor(0xF0, 0xF4, 0xF5), C_MUTED),
    ]
    for tit, url, desc, x, bg, border in puertas:
        left = Inches(x)
        add_round_box(slide, left, Inches(2.55), Inches(2.35), Inches(1.55), bg, border, 2)
        add_textbox(slide, left + Inches(0.1), Inches(2.65), Inches(2.15), Inches(0.35), tit, size=11, bold=True, color=C_ACCENT_DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), Inches(2.95), Inches(2.15), Inches(0.3), url, size=9, color=C_MUTED, align=PP_ALIGN.CENTER)
        add_textbox(slide, left + Inches(0.1), Inches(3.25), Inches(2.15), Inches(0.75), desc, size=9, color=C_TEXT, align=PP_ALIGN.CENTER)
        add_arrow(slide, left + Inches(0.95), Inches(2.35), Inches(0.35), Inches(0.22))

    box = add_round_box(slide, Inches(0.75), Inches(4.55), Inches(12), Inches(1.55), C_WHITE, C_WARN, 2)
    add_textbox(slide, Inches(0.95), Inches(4.7), Inches(11.6), Inches(1.25),
                "En palabras simples:\n"
                "• «Parcial, no gate del panel» = mova_auth/login.php existe, pero el panel /mova/ NO lo usa como entrada.\n"
                "• «Login aparte» = el ERP tiene su propia pantalla de clave, distinta de todo lo demás.\n"
                "• Conclusión D1: ningún módulo M tiene a mova_auth como único portero.",
                size=13, color=C_TEXT)


def slide_que_es_mova_auth(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "¿Qué es mova_auth?", "Glosario para no técnicos")

    cards = [
        ("mova_auth", "Carpeta en el servidor: public_html/acme-chile.cl/mova_auth/\nEs el «portero» que debería validar a todos antes de entrar a apps MOVA."),
        ("login.php", "Pantalla donde el usuario escribe correo y clave (o Google en el diseño futuro).\nHoy existe pero NO controla el acceso al panel /mova/."),
        ("guard.php", "Archivo que cada módulo privado incluirá al inicio (futuro).\nPregunta: «¿hay sesión?» — si no, manda al login."),
        ("Módulo M", "Cualquier app interna MOVA: portal, ERP, AXON, RRHH, facturas bajo /mova/, etc."),
    ]
    x = Inches(0.65)
    for tit, desc in cards:
        add_round_box(slide, x, Inches(1.35), Inches(2.95), Inches(2.35), C_WHITE, C_ACCENT)
        add_round_box(slide, x + Inches(0.1), Inches(1.48), Inches(2.75), Inches(0.42), C_ACCENT)
        add_textbox(slide, x + Inches(0.1), Inches(1.55), Inches(2.75), Inches(0.3), tit, size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(2.05), Inches(2.65), Inches(1.55), desc, size=10, color=C_TEXT)
        x += Inches(3.15)

    add_textbox(slide, Inches(0.75), Inches(4.0), Inches(12), Inches(0.35), "Analogía del edificio", size=14, bold=True, color=C_ACCENT_DARK)
    add_round_box(slide, Inches(0.75), Inches(4.4), Inches(12), Inches(1.85), C_PURPLE)
    add_textbox(slide, Inches(1.0), Inches(4.55), Inches(11.5), Inches(1.55),
                "Hoy: cada oficina (módulo) tiene su propia cerradura.\n"
                "Objetivo: un solo recepcionista (mova_auth) en la entrada. Si no pasaste por recepción, ninguna oficina te abre.",
                size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_regla_oro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Regla de oro", "Propuesta para acuerdo del equipo")
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.5), Inches(11.8), Inches(1.35))
    box.fill.solid()
    box.fill.fore_color.rgb = C_HIGHLIGHT
    box.line.color.rgb = C_WARN
    add_textbox(slide, Inches(0.95), Inches(1.7), Inches(11.4), Inches(1.0),
                "Si el usuario no pasó por mova_auth con sesión válida,\nno entra a ningún módulo M.",
                size=22, bold=True, color=C_ACCENT_DARK, align=PP_ALIGN.CENTER)
    pasos = [
        "1. Todo módulo privado incluye guard.php al inicio.",
        "2. Sesión PHP server-side — sin JWT en localStorage.",
        "3. Sin sesión → redirect a /mova_auth/login.php?redirect=…",
        "4. Tras login → vuelve al módulo pedido.",
        "5. Sin logins paralelos por módulo (salvo excepciones documentadas).",
    ]
    y = Inches(3.2)
    for p in pasos:
        add_textbox(slide, Inches(0.85), y, Inches(11.5), Inches(0.45), p, size=16, color=C_TEXT)
        y += Inches(0.5)
    gate = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(5.85), Inches(4.3), Inches(0.65))
    gate.fill.solid()
    gate.fill.fore_color.rgb = C_PURPLE
    gate.line.fill.background()
    add_textbox(slide, Inches(4.5), Inches(6.0), Inches(4.3), Inches(0.4), "mova_auth = única puerta", size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_regla_explicada(prs):
    """Cada punto de la regla de oro en lenguaje cotidiano."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Regla de oro — explicada paso a paso", "Qué significa cada punto en la práctica")

    pasos = [
        ("1 · guard.php", "Como un guardia en la puerta de cada módulo.\nAntes de mostrar contenido, verifica: «¿ya te autenticaste en mova_auth?»"),
        ("2 · Sesión en servidor", "La «credencial» vive en el servidor PHP, no en el navegador.\nNo guardamos tokens en localStorage (más seguro)."),
        ("3 · Redirect al login", "Si no hay sesión → pantalla login.php con parámetro redirect.\nEj: querías /mova/erp/ → tras login vuelves ahí."),
        ("4 · Vuelta al módulo", "El usuario no pierde su destino. Entra una vez y llega donde iba."),
        ("5 · Sin logins paralelos", "Se acaban las pantallas de Google sueltas, claves ERP aparte, etc.\nSolo mova_auth (salvo excepciones escritas)."),
    ]
    y = Inches(1.35)
    for tit, desc in pasos:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), y + Inches(0.08), Inches(0.35), Inches(0.35))
        circ.fill.solid()
        circ.fill.fore_color.rgb = C_PURPLE
        circ.line.fill.background()
        add_round_box(slide, Inches(1.15), y, Inches(11.5), Inches(0.95), C_WHITE, C_PURPLE)
        add_textbox(slide, Inches(1.3), y + Inches(0.06), Inches(2.5), Inches(0.28), tit, size=11, bold=True, color=C_PURPLE)
        add_textbox(slide, Inches(1.3), y + Inches(0.34), Inches(11.1), Inches(0.55), desc, size=10, color=C_TEXT)
        y += Inches(1.05)


def slide_flujo_objetivo(prs):
    """Mapa conceptual del flujo futuro deseado."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Mapa conceptual — flujo OBJETIVO", "Cómo debería funcionar (diseño acordado)")

    steps = [
        ("1", "Usuario pide\n/mova/erp/", C_MUTED),
        ("2", "guard.php\n¿sesión?", C_ACCENT),
        ("3", "No → login.php\n?redirect=/mova/erp/", C_PURPLE),
        ("4", "Login OK\nsesión + cookie", C_OK),
        ("5", "Vuelve a\n/mova/erp/", C_ACCENT),
    ]
    x = Inches(0.55)
    for num, label, col in steps:
        add_round_box(slide, x, Inches(2.0), Inches(2.2), Inches(1.15), C_WHITE, col, 2)
        add_textbox(slide, x, Inches(2.05), Inches(2.2), Inches(0.3), num, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.1), Inches(2.4), Inches(2.0), Inches(0.65), label, size=10, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
        if x < Inches(10):
            add_arrow(slide, x + Inches(2.25), Inches(2.45), Inches(0.35), Inches(0.25), col)
        x += Inches(2.55)

    add_round_box(slide, Inches(0.75), Inches(3.65), Inches(5.8), Inches(2.35), C_GREEN_BG, C_OK, 2)
    add_textbox(slide, Inches(0.95), Inches(3.8), Inches(5.4), Inches(0.3), "Con sesión válida", size=12, bold=True, color=C_OK)
    add_textbox(slide, Inches(0.95), Inches(4.15), Inches(5.4), Inches(1.7),
                "guard.php deja pasar.\nEl módulo carga normal.\nLa cookie HttpOnly viaja sola\n(sin que el usuario haga nada).", size=11, color=C_TEXT)

    add_round_box(slide, Inches(6.85), Inches(3.65), Inches(5.8), Inches(2.35), C_ORANGE_BG, C_WARN, 2)
    add_textbox(slide, Inches(7.05), Inches(3.8), Inches(5.4), Inches(0.3), "Sin sesión", size=12, bold=True, color=C_WARN)
    add_textbox(slide, Inches(7.05), Inches(4.15), Inches(5.4), Inches(1.7),
                "guard.php redirige a login.\nEl usuario NO ve contenido privado.\nTras autenticarse, regresa\nal módulo que pidió.", size=11, color=C_TEXT)

    add_textbox(slide, Inches(0.75), Inches(6.25), Inches(12), Inches(0.4),
                "Clave: el módulo nunca decide solo quién entra — solo pregunta a mova_auth vía guard.php.",
                size=13, bold=True, color=C_ACCENT_DARK)


def slide_tabla(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "¿Quién debe pasar por mova_auth?", "Basado en inventario D1")
    headers = ["Módulo", "URL", "Auth hoy", "¿mova_auth?", "Notas"]
    xs = [Inches(0.6), Inches(2.4), Inches(5.2), Inches(8.0), Inches(9.8)]
    for i, h in enumerate(headers):
        add_textbox(slide, xs[i], Inches(1.35), Inches(2.2), Inches(0.35), h, size=11, bold=True, color=C_ACCENT_DARK)
    y = Inches(1.75)
    for mod, url, auth, mova, nota in MODULOS:
        col = C_OK if mova == "Sí" else C_PURPLE if "gate" in mova else C_WARN if mova == "No" else C_TEXT
        add_textbox(slide, xs[0], y, Inches(1.7), Inches(0.32), mod, size=10, bold=True, color=C_TEXT)
        add_textbox(slide, xs[1], y, Inches(2.7), Inches(0.32), url, size=9, color=C_MUTED)
        add_textbox(slide, xs[2], y, Inches(2.6), Inches(0.32), auth, size=9, color=C_TEXT)
        add_textbox(slide, xs[3], y, Inches(1.6), Inches(0.32), mova, size=10, bold=True, color=col)
        add_textbox(slide, xs[4], y, Inches(2.8), Inches(0.32), nota, size=9, color=C_TEXT)
        y += Inches(0.4)
    add_textbox(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.5),
                "Submódulos bajo /mova/ (facturas, oc, agencia…) heredan la misma regla: deben pasar por mova_auth.",
                size=13, color=C_ACCENT_DARK)


def slide_tabla_explicada(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Tabla explicada — fila por fila", "Para compartir con quien no vio el inventario D1")

    filas = [
        ("Portal MOVA", "Hoy entra con Google. En el diseño nuevo debe pasar por mova_auth primero."),
        ("mova_auth", "Es el portero. Hoy es «parcial» porque existe pero no controla todo."),
        ("MOVA ERP", "Hoy tiene clave propia. Debe unificarse: misma sesión que el resto."),
        ("AXON", "Abre sin login visible. Hay que definir validación (tarea D5)."),
        ("RRHH", "403 en servidor. Privado — debe pasar por mova_auth cuando se habilite."),
        ("Documentos", "Público a propósito. EXCEPCIÓN — no necesita mova_auth."),
    ]
    y = Inches(1.3)
    for tit, desc in filas:
        add_round_box(slide, Inches(0.65), y, Inches(2.2), Inches(0.72), C_ACCENT)
        add_textbox(slide, Inches(0.65), y + Inches(0.18), Inches(2.2), Inches(0.4), tit, size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_round_box(slide, Inches(3.0), y, Inches(9.65), Inches(0.72), C_WHITE, C_MUTED)
        add_textbox(slide, Inches(3.15), y + Inches(0.15), Inches(9.35), Inches(0.45), desc, size=11, color=C_TEXT)
        y += Inches(0.82)

    add_textbox(slide, Inches(0.65), Inches(6.35), Inches(12), Inches(0.35),
                "Lectura rápida: verde=Sí debe pasar · violeta=es el gate · naranja=excepción pública",
                size=11, bold=True, color=C_ACCENT_DARK)


def slide_excepciones(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Excepciones y contenido público", "No requieren mova_auth")
    items = [
        ("/documentos/", "Playbooks HTML públicos (auditoria_mova.html, etc.)"),
        ("Assets estáticos", "CSS, JS, imágenes sin datos sensibles"),
        ("login.php", "Página de entrada — no se protege con guard"),
        ("logout.php", "Cierre de sesión"),
        ("validate.php", "API JSON 200/401 para AJAX"),
    ]
    y = Inches(1.55)
    for tit, desc in items:
        add_textbox(slide, Inches(0.85), y, Inches(3.5), Inches(0.35), tit, size=14, bold=True, color=C_ACCENT_DARK)
        add_textbox(slide, Inches(4.2), y, Inches(8.2), Inches(0.35), desc, size=13, color=C_TEXT)
        y += Inches(0.55)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.8), Inches(11.8), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = C_ORANGE_BG
    box.line.color.rgb = C_WARN
    add_textbox(slide, Inches(0.95), Inches(5.0), Inches(11.4), Inches(0.7),
                "Prohibido post-acuerdo: Google OAuth directo en módulo · JWT en localStorage · logins duplicados por módulo.",
                size=14, bold=True, color=C_WARN)


def slide_excepciones_mapa(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Mapa — ¿Qué está adentro y afuera?", "Dentro del edificio vs calle pública")

    add_round_box(slide, Inches(0.75), Inches(1.35), Inches(5.5), Inches(4.85), C_GREEN_BG, C_OK, 3)
    add_textbox(slide, Inches(1.0), Inches(1.5), Inches(5.0), Inches(0.35), "DENTRO — requiere mova_auth", size=14, bold=True, color=C_OK)
    for i, mod in enumerate(["/mova/ y submódulos", "/mova/erp/", "/axon/", "/rrhh/", "Todo módulo M privado"]):
        add_round_box(slide, Inches(1.1), Inches(2.0 + i * 0.65), Inches(4.8), Inches(0.5), C_WHITE, C_OK)
        add_textbox(slide, Inches(1.1), Inches(2.12 + i * 0.65), Inches(4.8), Inches(0.3), mod, size=11, color=C_TEXT, align=PP_ALIGN.CENTER)

    add_round_box(slide, Inches(6.95), Inches(1.35), Inches(5.65), Inches(4.85), RGBColor(0xF0, 0xF4, 0xF5), C_MUTED, 3)
    add_textbox(slide, Inches(7.2), Inches(1.5), Inches(5.2), Inches(0.35), "AFUERA — público / excepción", size=14, bold=True, color=C_MUTED)
    for i, mod in enumerate(["/documentos/ playbooks", "CSS · JS · imágenes", "login.php (entrada)", "logout.php", "validate.php API"]):
        add_round_box(slide, Inches(7.3), Inches(2.0 + i * 0.65), Inches(4.95), Inches(0.5), C_WHITE, C_MUTED)
        add_textbox(slide, Inches(7.3), Inches(2.12 + i * 0.65), Inches(4.95), Inches(0.3), mod, size=11, color=C_TEXT, align=PP_ALIGN.CENTER)

    wall = add_round_box(slide, Inches(6.35), Inches(2.5), Inches(0.55), Inches(2.5), C_PURPLE)
    add_textbox(slide, Inches(6.35), Inches(3.2), Inches(0.55), Inches(1.0), "mova\n_auth", size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.75), Inches(6.35), Inches(12), Inches(0.4),
                "Prohibido: poner Google OAuth o JWT dentro de un módulo «adentro» sin pasar por el portero.",
                size=12, bold=True, color=C_WARN)


def slide_antes_despues(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Antes vs después (resumen visual)", "Lo que acordamos cambiar")

    add_textbox(slide, Inches(0.85), Inches(1.35), Inches(5.5), Inches(0.35), "ANTES (Día 1)", size=16, bold=True, color=C_WARN)
    antes = ["Google en /mova/", "Clave suelta en ERP", "mova_auth/login.php aislado", "Sin portero único"]
    y = Inches(1.8)
    for a in antes:
        add_round_box(slide, Inches(0.85), y, Inches(5.5), Inches(0.48), C_ORANGE_BG, C_WARN)
        add_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(5.2), Inches(0.3), a, size=12, color=C_WARN, align=PP_ALIGN.CENTER)
        y += Inches(0.55)

    add_arrow(slide, Inches(6.55), Inches(3.2), Inches(0.6), Inches(0.4), C_PURPLE)

    add_textbox(slide, Inches(7.35), Inches(1.35), Inches(5.5), Inches(0.35), "DESPUÉS (acuerdo D2)", size=16, bold=True, color=C_OK)
    despues = ["Un solo login en mova_auth", "Sesión PHP + cookie HttpOnly", "guard.php en cada módulo", "Sin JWT en navegador"]
    y = Inches(1.8)
    for d in despues:
        add_round_box(slide, Inches(7.35), y, Inches(5.5), Inches(0.48), C_GREEN_BG, C_OK)
        add_textbox(slide, Inches(7.5), y + Inches(0.1), Inches(5.2), Inches(0.3), d, size=12, color=C_OK, align=PP_ALIGN.CENTER)
        y += Inches(0.55)

    add_round_box(slide, Inches(0.75), Inches(4.85), Inches(12), Inches(1.35), C_HIGHLIGHT, C_ACCENT)
    add_textbox(slide, Inches(0.95), Inches(5.05), Inches(11.6), Inches(1.0),
                "Importante: «después» es el diseño acordado en auditoría.\n"
                "La implementación en servidor empieza en Día 3+. Hoy solo firmamos las reglas.",
                size=14, bold=True, color=C_ACCENT_DARK, align=PP_ALIGN.CENTER)


def slide_checklist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Checklist cierre Día 2", "Compartir con equipo técnico")
    y = Inches(1.5)
    for label, ok in CHECKLIST:
        mark = "✓" if ok else "○"
        col = C_OK if ok else C_WARN
        add_textbox(slide, Inches(0.85), y, Inches(0.5), Inches(0.4), mark, size=20, bold=True, color=col)
        add_textbox(slide, Inches(1.35), y, Inches(10.5), Inches(0.4), label, size=18, color=C_TEXT)
        y += Inches(0.52)
    add_textbox(slide, Inches(0.85), Inches(5.5), Inches(11.5), Inches(0.5),
                "Pendiente: acuerdo formal (correo/acta) · Responsable: _______ · Fecha: _______", size=14, color=C_MUTED)
    add_textbox(slide, Inches(0.85), Inches(6.1), Inches(11.5), Inches(0.45),
                "Próximo: Día 3 — Carpetas y archivos núcleo · mkof/03", size=15, bold=True, color=C_ACCENT_DARK)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_title(prs)
    slide_contexto(prs)
    slide_puente(
        prs,
        "Hallazgos D1",
        "Mapa situación HOY",
        "El inventario listó cuatro formas distintas de «entrar». Antes de proponer reglas, "
        "visualizamos el problema: el usuario no sabe qué puerta usar y ningún módulo comparte la misma sesión.",
        [
            "Inventario D1  →  «login fragmentado»  →  4 puertas distintas",
            "  /mova/ (Google)     mova_auth/login.php     /mova/erp/ (clave)     /axon/ (¿?) ",
            "Conclusión: sin portero único, no hay seguridad ni experiencia coherente.",
        ],
    )
    slide_mapa_fragmentado(prs)
    slide_puente(
        prs,
        "Mapa HOY",
        "¿Qué es mova_auth?",
        "Vimos el caos actual. Ahora definimos el vocabulario del acuerdo: qué es mova_auth, "
        "qué hace guard.php y qué entendemos por «módulo M» — para que todos hablen lo mismo.",
        [
            "Problema (varias puertas)  →  Solución propuesta: UN recepcionista",
            "mova_auth = carpeta portero  |  guard.php = guardia en cada módulo  |  M = app interna MOVA",
        ],
    )
    slide_que_es_mova_auth(prs)
    slide_puente(
        prs,
        "Glosario",
        "Regla de oro",
        "Con los términos claros, pasamos a la regla que el equipo debe acordar por escrito. "
        "Es la «constitución» del acceso: una sola frase que decide quién entra y quién no.",
        [
            "Analogía edificio  →  Regla escrita  →  5 reglas operativas numeradas",
            "Si no pasaste por mova_auth con sesión válida → ningún módulo M te abre.",
        ],
    )
    slide_regla_oro(prs)
    slide_regla_explicada(prs)
    slide_puente(
        prs,
        "Regla de oro",
        "Flujo OBJETIVO",
        "La regla en papel se traduce en un recorrido concreto. Este mapa muestra qué pasa "
        "cuando un usuario pide un módulo: guard verifica, login si hace falta, y vuelta al destino.",
        [
            "Usuario → módulo → guard.php → ¿sesión? → Sí: entra  |  No: login.php?redirect= → vuelve",
        ],
    )
    slide_flujo_objetivo(prs)
    slide_puente(
        prs,
        "Flujo objetivo",
        "Tabla de módulos",
        "El flujo aplica a apps concretas. La tabla del inventario D1 dice, módulo por módulo, "
        "quién debe pasar por mova_auth y cuál es la excepción documentada.",
        [
            "Diseño genérico (guard + login)  →  Inventario real: Portal, ERP, AXON, RRHH, Documentos…",
        ],
    )
    slide_tabla(prs)
    slide_tabla_explicada(prs)
    slide_puente(
        prs,
        "Tabla módulos",
        "Excepciones públicas",
        "No todo el sitio es privado. Hay rutas que deben quedar abiertas (playbooks, assets, "
        "login mismo). También listamos lo prohibido después del acuerdo.",
        [
            "DENTRO (módulos M)  ← mova_auth →  AFUERA (público / API validate)",
            "Prohibido: Google suelto en módulo · JWT en localStorage · logins duplicados",
        ],
    )
    slide_excepciones(prs)
    slide_excepciones_mapa(prs)
    slide_puente(
        prs,
        "Excepciones",
        "Antes vs después",
        "Cerramos con un resumen visual del cambio acordado. Importante: hoy solo documentamos; "
        "la implementación en servidor empieza en Día 3+.",
        [
            "ANTES: 4 logins · sesiones sueltas     →     DESPUÉS: 1 login · guard.php · cookie HttpOnly",
        ],
    )
    slide_antes_despues(prs)
    slide_checklist(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"PPT generado: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Tamaño: {OUT.stat().st_size // 1024} KB")


if __name__ == "__main__":
    main()
