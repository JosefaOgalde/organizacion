#!/usr/bin/env python3
"""Genera MOVA-Nodo-MySQL-n8n.pptx — paso a paso agregar nodo MySQL en n8n (hito 3.2)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

OUT = Path(__file__).resolve().parents[1] / "index/clientes/mkof/MOVA-Nodo-MySQL-n8n.pptx"

# Dimensiones / colores (igual que guías GitHub MOVA)
W = 12191695
H = 6858000
TEAL = RGBColor(0x4A, 0x7A, 0x80)
TEAL_DARK = RGBColor(0x2A, 0x4A, 0x4E)
INK = RGBColor(0x1F, 0x23, 0x28)
MUTED = RGBColor(0x65, 0x6D, 0x76)
BLUE = RGBColor(0x09, 0x69, 0xDA)
GOLD = RGBColor(0xD4, 0xA7, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG = RGBColor(0xF6, 0xF8, 0xFA)
YELLOW = RGBColor(0xFF, 0xF8, 0xC5)
MINT = RGBColor(0xE8, 0xF6, 0xF8)
BORDER = RGBColor(0xD0, 0xD7, 0xDE)

PASOS = [
    {
        "titulo": "Abrir n8n y el workflow",
        "texto": "Entra a la instancia n8n de MOVA y abre el workflow que hoy escribe a Google Sheets (ingresos/egresos).",
        "tip": "Si hay varios workflows, elige el que alimenta MOVA MAESTRO / Sheets de operación.",
        "mock_url": "n8n → Workflows",
        "mock_lines": ["Workflows", "MOVA · Ingresos / Egresos", "Estado: Active · Editor"],
        "highlight": "👆 Abrir el workflow que hoy escribe a Sheets",
    },
    {
        "titulo": "Agregar un nodo nuevo",
        "texto": "En el canvas, haz clic en el botón ＋ (Add node) al final de la cadena o en el punto de conexión donde irá MySQL.",
        "tip": "MySQL debe quedar ANTES de Google Sheets: primero escribe a BD, después replica a la vista.",
        "mock_url": "n8n → canvas",
        "mock_lines": ["Trigger / Webhook", "→  ＋ Add node", "Google Sheets (hoy)"],
        "highlight": "👆 Botón ＋ Add node",
    },
    {
        "titulo": "Buscar el nodo MySQL",
        "texto": "En el buscador de nodos escribe MySQL y selecciona el nodo MySQL (no Postgres ni Microsoft SQL).",
        "tip": "El nodo oficial se llama exactamente «MySQL». Evita nodos community no aprobados por el equipo.",
        "mock_url": "n8n → Add node",
        "mock_lines": ["Search nodes", "MySQL  ← AQUÍ", "MySQL · Insert / Update / Select"],
        "highlight": "👆 Escribe MySQL y elige el nodo oficial",
    },
    {
        "titulo": "Crear credencial MySQL",
        "texto": "En Credential → Create New. Nombre sugerido: MOVA MySQL prod (o sandbox si estás en prueba).",
        "tip": "Una sola credencial compartida por workflows. No pegues la contraseña en el chat ni en el repo Git.",
        "mock_url": "n8n → MySQL → Credential",
        "mock_lines": ["Credential", "Create New", "Name: MOVA MySQL prod"],
        "highlight": "👆 Create New → nombre de credencial del equipo",
    },
    {
        "titulo": "Completar Host, Database, User, Password",
        "texto": "Usa los datos del hito 3.1 (MySQL gestionada). Host, Port 3306, Database mova_datos, User con permisos mínimos, Password del gestor.",
        "tip": "Si n8n Cloud no alcanza la BD: revisa whitelist de IP (hito 3.1). Sin IP autorizada, Test fallará.",
        "mock_url": "n8n → MySQL credential",
        "mock_lines": ["Host: xxx.mysql.cloud", "Database: mova_datos", "User / Password · Port 3306"],
        "highlight": "👆 Host · Database · User · Password (hito 3.1)",
    },
    {
        "titulo": "Probar la conexión (Test)",
        "texto": "Guarda la credencial y pulsa Test / Retry. Debe aparecer Connection successful antes de seguir.",
        "tip": "Si falla: IP whitelist, SSL requerido, usuario/clave o firewall. No configures el query hasta que el test pase.",
        "mock_url": "n8n → Credential → Test",
        "mock_lines": ["Test connection", "✓ Connection successful", "Save"],
        "highlight": "👆 Test → Connection successful",
    },
    {
        "titulo": "Elegir la operación SQL",
        "texto": "Operation típica para escritura: Insert o Insert or Update. Para lecturas/reportes: Execute SQL / Select.",
        "tip": "Para la fuente primaria de MOVA: Insert or Update con clave única (ej. id_movimiento) evita duplicados.",
        "mock_url": "n8n → MySQL node",
        "mock_lines": ["Operation", "Insert or Update  ← AQUÍ", "Table: ingresos"],
        "highlight": "👆 Insert or Update + tabla del esquema 3.1",
    },
    {
        "titulo": "Mapear columnas desde el nodo anterior",
        "texto": "En Columns / Values, mapea cada campo desde el JSON del nodo previo (fecha, monto, tipo, módulo, etc.).",
        "tip": "Usa expresiones {{ $json.campo }}. No hardcodees montos. Revisa tipos DATE / DECIMAL en el esquema.",
        "mock_url": "n8n → MySQL → Columns",
        "mock_lines": ["fecha → {{ $json.fecha }}", "monto → {{ $json.monto }}", "tipo → {{ $json.tipo }}"],
        "highlight": "👆 Mapear columnas 1:1 con el esquema mova_datos",
    },
    {
        "titulo": "Orden del flujo: MySQL → luego Sheets",
        "texto": "Conecta: … → MySQL → Google Sheets. MySQL es fuente primaria; Sheets queda como vista de lectura.",
        "tip": "Si Sheets falla y MySQL ya escribió, los datos siguen seguros en BD. Al revés, no.",
        "mock_url": "n8n → canvas",
        "mock_lines": ["Webhook / lógica", "→ MySQL (fuente)", "→ Google Sheets (vista)"],
        "highlight": "👆 MySQL antes de Sheets en el canvas",
    },
    {
        "titulo": "Configurar Sheets como réplica",
        "texto": "En el nodo Google Sheets: Append o Update con los mismos campos. Solo lectura operativa para el equipo.",
        "tip": "No uses Sheets como fuente de verdad. Si alguien edita la hoja a mano, la BD manda.",
        "mock_url": "n8n → Google Sheets",
        "mock_lines": ["Operation: Append", "Sheet: INGRESOS / EGRESOS", "Mismos campos que MySQL"],
        "highlight": "👆 Sheets = vista · MySQL = verdad",
    },
    {
        "titulo": "Ejecutar una prueba",
        "texto": "Execute Workflow (o Test step en MySQL). Verifica fila nueva en MySQL y réplica en Sheets.",
        "tip": "Prueba primero en sandbox / tabla de prueba. Luego activa el workflow en producción.",
        "mock_url": "n8n → Executions",
        "mock_lines": ["Execute Workflow", "MySQL · Success", "Sheets · Success"],
        "highlight": "👆 Success en MySQL y en Sheets",
    },
]


def set_run(paragraph, text, *, size=18, bold=False, color=INK, name=None):
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if name:
        run.font.name = name
    return run


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p, text, size=size, bold=bold, color=color)
    return box


def add_rect(slide, left, top, width, height, fill, *, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, Emu(left), Emu(top), Emu(width), Emu(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def add_oval(slide, left, top, size, fill, text):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(left), Emu(top), Emu(size), Emu(size))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    tf = s.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, text, size=18, bold=True, color=WHITE)
    return s


def header_bar(slide, title, subtitle):
    add_rect(slide, 0, 0, W, 960120, TEAL)
    add_textbox(slide, 502920, 164592, 9144000, 457200, title, size=22, bold=True, color=WHITE)
    add_textbox(slide, 502920, 566928, 9144000, 320040, subtitle, size=14, color=MINT)


def slide_portada(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 731520, 1828800, 10515600, 1097280, "MOVA · n8n + MySQL", size=36, bold=True, color=TEAL_DARK)
    add_textbox(slide, 731520, 2926080, 10515600, 731520, "Paso a paso: agregar el nodo MySQL", size=24, color=INK)
    add_textbox(
        slide,
        731520,
        3931920,
        10515600,
        457200,
        "Hito 3.2 · MySQL fuente → Sheets vista  |  GRUPO MAKING OF",
        size=14,
        color=MUTED,
    )
    add_textbox(slide, 731520, 5669280, 10515600, 365760, "Guía para el encargado · Julio 2026", size=13, color=MUTED)


def slide_objetivo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "¿Qué vamos a hacer?", "Objetivo de esta guía")
    bullets = [
        "•  Agregar el nodo MySQL en el workflow n8n de MOVA.",
        "•  Configurar la credencial con la base mova_datos (hito 3.1).",
        "•  Escribir primero a MySQL (fuente primaria).",
        "•  Replicar después a Google Sheets (solo vista de lectura).",
        "•  Tiempo estimado: 20–30 minutos (con BD ya provisionada).",
        "•  Prerrequisito: hito 3.1 — esquema MySQL + whitelist IP listos.",
    ]
    y = 1371600
    for b in bullets:
        add_textbox(slide, 822960, y, 10058400, 502920, b, size=16, color=INK)
        y += 594360


def slide_requisitos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Antes de empezar", "Hito 3.1 cerrado · datos a mano")
    tags = ["Host MySQL", "DB mova_datos", "User + pass", "IP whitelist"]
    x = 640080
    for t in tags:
        add_rect(slide, x, 1325880, 2377440, 411480, WHITE, rounded=True)
        add_textbox(slide, x + 91440, 1389888, 2194560, 320040, t, size=12, bold=True, color=TEAL_DARK)
        x += 2514600
    checks = [
        "✓  MySQL gestionada creada (cloud) con esquema ingresos/egresos.",
        "✓  IP de n8n (o rango) en whitelist de la BD.",
        "✓  Usuario con permisos mínimos (INSERT/UPDATE/SELECT, no root).",
        "✓  Credenciales guardadas en gestor del equipo (1Password / Bitwarden).",
    ]
    y = 2103120
    for c in checks:
        add_textbox(slide, 822960, y, 10058400, 457200, c, size=15, color=INK)
        y += 502920
    add_rect(slide, 640080, 4389120, 10515600, 1463040, YELLOW, rounded=True)
    add_textbox(
        slide,
        868680,
        4572000,
        10058400,
        1097280,
        "⚠  Sin hito 3.1 listo, el Test de conexión fallará. No inventes host/clave: "
        "pídelos al responsable de infraestructura MOVA.",
        size=14,
        color=INK,
    )


def slide_paso(prs, num: int, total: int, paso: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, f"Paso {num} de {total}", paso["titulo"])

    add_oval(slide, 502920, 1234440, 502920, TEAL, str(num))
    add_textbox(slide, 1143000, 1143000, 4754880, 502920, paso["titulo"], size=18, bold=True, color=INK)
    add_textbox(slide, 685800, 1783080, 5029200, 1645920, paso["texto"], size=15, color=INK)

    add_rect(slide, 640080, 3657600, 5120640, 1371600, GRAY_BG, rounded=True)
    add_textbox(slide, 822960, 3794760, 4754880, 1097280, f"💡 Tip: {paso['tip']}", size=13, color=MUTED)

    # Mock panel derecha
    add_rect(slide, 6217920, 1371600, 5303520, 4754880, WHITE, rounded=True)
    add_rect(slide, 6400800, 1554480, 4937760, 411480, GRAY_BG, rounded=True)
    add_textbox(slide, 6537960, 1627632, 4754880, 274320, paso["mock_url"], size=12, bold=True, color=TEAL_DARK)
    add_textbox(slide, 6537960, 2148840, 4663440, 365760, "n8n · MySQL", size=16, bold=True, color=INK)

    y = 2606040
    for i, line in enumerate(paso["mock_lines"]):
        add_textbox(slide, 6537960, y, 4663440, 320040, line, size=13, color=INK if i else MUTED)
        y += 365760

    add_rect(slide, 6537960, 4800600, 4663440, 820000, YELLOW, rounded=True)
    add_textbox(slide, 6675120, 4937760, 4389120, 640080, paso["highlight"], size=12, color=INK)


def slide_checklist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, "Checklist final", "Marca cada ítem antes de avisar que terminaste")
    items = [
        "Credencial MySQL creada y Test OK",
        "Nodo MySQL en el workflow (Insert or Update)",
        "Columnas mapeadas al esquema mova_datos",
        "Orden canvas: MySQL → Google Sheets",
        "Ejecución de prueba Success en ambos nodos",
        "Workflow Active en producción (o sandbox validado)",
    ]
    y = 1371600
    for item in items:
        add_rect(slide, 731520, y, 320040, 320040, WHITE, rounded=True)
        add_textbox(slide, 1234440, y - 45720, 9601200, 411480, item, size=15, color=INK)
        y += 594360


def slide_cierre(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 731520, 2011680, 10515600, 640080, "Hito 3.2 · listo para operación", size=22, bold=True, color=TEAL_DARK)
    box = slide.shapes.add_textbox(Emu(731520), Emu(2926080), Emu(10515600), Emu(1097280))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "Siguiente: "
    r1.font.size = Pt(24)
    r1.font.color.rgb = INK
    r2 = p.add_run()
    r2.text = "Rutinas de operación (hito 4.0)"
    r2.font.size = Pt(24)
    r2.font.bold = True
    r2.font.color.rgb = TEAL

    add_textbox(
        slide,
        731520,
        4572000,
        10515600,
        731520,
        "Guía: mysql-nodo-n8n.html · PPT: MOVA-Nodo-MySQL-n8n.pptx",
        size=14,
        color=MUTED,
    )
    add_textbox(slide, 731520, 5760720, 10515600, 365760, "MOVA · GRUPO MAKING OF · Jul 2026", size=13, color=MUTED)


def main():
    prs = Presentation()
    prs.slide_width = Emu(W)
    prs.slide_height = Emu(H)

    slide_portada(prs)
    slide_objetivo(prs)
    slide_requisitos(prs)
    total = len(PASOS)
    for i, paso in enumerate(PASOS, start=1):
        slide_paso(prs, i, total, paso)
    slide_checklist(prs)
    slide_cierre(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"OK → {OUT} ({2 + 1 + total + 2} slides)")


if __name__ == "__main__":
    main()
