#!/usr/bin/env python3
"""Genera PPTX · Talento TP Digital para Chile · estética desafiolatam.com"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from copy import deepcopy

OUT = Path(__file__).resolve().parents[1] / "index/clientes/DesafioLatam/Talento-TP-Digital/Talento-TP-Digital-Banco-de-Chile.pptx"

INK = RGBColor(0x0F, 0x17, 0x22)
GREEN = RGBColor(0x5D, 0x8A, 0x22)
GREEN_SOFT = RGBColor(0xE4, 0xEF, 0xC7)
GRAY = RGBColor(0xF2, 0xF4, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x4A, 0x55, 0x63)

W, H = Inches(13.333), Inches(7.5)


def set_run(run, size=18, bold=False, color=INK, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_text(shape, lines, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, italic=False):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(line, tuple):
            text, *opts = line
            r = p.add_run()
            r.text = text
            set_run(
                r,
                size=opts[0] if len(opts) > 0 else size,
                bold=opts[1] if len(opts) > 1 else bold,
                color=opts[2] if len(opts) > 2 else color,
                italic=opts[3] if len(opts) > 3 else italic,
            )
        else:
            r = p.add_run()
            r.text = line
            set_run(r, size=size, bold=bold, color=color, italic=italic)
    return tf


def rect(slide, left, top, width, height, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.adjustments[0] = 0.1
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill or WHITE
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def pill(slide, left, top, width, height, text, fill=GREEN, font_size=18):
    sh = rect(slide, left, top, width, height, fill=fill)
    sh.adjustments[0] = 0.35
    tf = sh.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    set_run(r, size=font_size, bold=True, color=WHITE)
    tf.paragraphs[0].space_before = Pt(0)
    try:
        tf.auto_size = None
    except Exception:
        pass
    return sh


def eyebrow(slide, text, left=Inches(0.7), top=Inches(0.45)):
    box = slide.shapes.add_textbox(left, top, Inches(10), Inches(0.35))
    add_text(box, [text], size=12, bold=True, color=GREEN)


def footer(slide, num):
    logo = slide.shapes.add_textbox(Inches(0.7), Inches(7.05), Inches(4), Inches(0.3))
    tf = logo.text_frame
    p = tf.paragraphs[0]
    r1 = p.add_run()
    r1.text = "{desafío}"
    set_run(r1, size=12, bold=True, color=GREEN)
    r2 = p.add_run()
    r2.text = " latam_"
    set_run(r2, size=12, bold=True, color=INK)
    n = slide.shapes.add_textbox(Inches(12.2), Inches(7.05), Inches(0.8), Inches(0.3))
    add_text(n, [f"{num:02d}"], size=12, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)


def title_block(slide, title_lines, pill_text=None, pill_w=Inches(2.4)):
    """title_lines without the pill part; pill_text goes in green pill below or inline"""
    y = Inches(0.9)
    for i, line in enumerate(title_lines):
        box = slide.shapes.add_textbox(Inches(0.7), y, Inches(7.2), Inches(0.85))
        add_text(box, [line], size=40, bold=True, color=INK)
        y += Inches(0.72)
    if pill_text:
        pill(slide, Inches(0.7), y, pill_w, Inches(0.55), pill_text, font_size=22)
        y += Inches(0.7)
    return y


def card(slide, left, top, width, height, title, body_lines, badge=None):
    sh = rect(slide, left, top, width, height, fill=GRAY, line=RGBColor(0xD4, 0xD8, 0xDD))
    # left accent
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    acc.fill.solid()
    acc.fill.fore_color.rgb = GREEN
    acc.line.fill.background()
    y = top + Inches(0.18)
    if badge:
        pill(slide, left + Inches(0.25), y, Inches(2.2), Inches(0.32), badge, font_size=10)
        y += Inches(0.42)
    tb = slide.shapes.add_textbox(left + Inches(0.25), y, width - Inches(0.4), Inches(0.4))
    add_text(tb, [title], size=16, bold=True, color=INK)
    y += Inches(0.4)
    body = slide.shapes.add_textbox(left + Inches(0.25), y, width - Inches(0.4), height - (y - top) - Inches(0.15))
    add_text(body, body_lines, size=13, color=MUTED)
    return sh


def metric(slide, left, top, width, height, num, label):
    sh = rect(slide, left, top, width, height, fill=WHITE, line=RGBColor(0xD4, 0xD8, 0xDD))
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.08), height)
    acc.fill.solid()
    acc.fill.fore_color.rgb = GREEN
    acc.line.fill.background()
    n = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.25), width - Inches(0.35), Inches(0.7))
    add_text(n, [num], size=32, bold=True, color=INK)
    l = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.95), width - Inches(0.35), Inches(0.5))
    add_text(l, [label.upper()], size=11, bold=True, color=GREEN)


def quote_box(slide, left, top, width, text):
    bg = rect(slide, left, top, width, Inches(1.15), fill=GREEN_SOFT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.1), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.25), width - Inches(0.45), Inches(0.8))
    add_text(tb, [text], size=15, bold=True, color=INK, italic=False)


def new_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 01 Portada
    s = new_slide(prs)
    eyebrow(s, "REPORTE DE IMPACTO · ABRIL–JULIO 2026")
    y = title_block(s, ["Talento TP Digital"], "para Chile", Inches(2.6))
    body = s.shapes.add_textbox(Inches(0.7), y + Inches(0.1), Inches(6.8), Inches(1.2))
    add_text(
        body,
        [
            "Capacitación en Inteligencia Artificial, Automatización y Competencias Financieras para jóvenes de Enseñanza Media TP."
        ],
        size=16,
        color=MUTED,
    )
    quote_box(
        s,
        Inches(0.7),
        Inches(4.55),
        Inches(6.8),
        "“El talento técnico-profesional no hay que salir a buscarlo: hay que prepararlo.”",
    )
    # right visual panel
    panel = rect(s, Inches(8.0), Inches(1.0), Inches(4.6), Inches(5.2), fill=INK)
    panel.adjustments[0] = 0.08
    dots = rect(s, Inches(10.2), Inches(0.7), Inches(2.6), Inches(2.4), fill=GREEN_SOFT)
    dots.adjustments[0] = 0.12
    floatc = rect(s, Inches(8.35), Inches(4.35), Inches(4.0), Inches(2.0), fill=WHITE)
    pill(s, Inches(8.55), Inches(4.5), Inches(1.6), Inches(0.32), "CASO DE ÉXITO", font_size=9)
    t = s.shapes.add_textbox(Inches(8.55), Inches(4.95), Inches(3.6), Inches(0.45))
    add_text(t, ["+85 estudiantes"], size=22, bold=True, color=INK)
    sub = s.shapes.add_textbox(Inches(8.55), Inches(5.4), Inches(3.6), Inches(0.35))
    add_text(sub, ["Initec Curicó · Región del Maule"], size=12, color=MUTED)
    for i, (n, lab) in enumerate([("92,5%", "Aprob."), ("202h", "Horas"), ("10/10", "Docentes")]):
        x = Inches(8.55) + Inches(i * 1.25)
        m = rect(s, x, Inches(5.85), Inches(1.15), Inches(0.7), fill=GRAY)
        tb = s.shapes.add_textbox(x + Inches(0.05), Inches(5.9), Inches(1.05), Inches(0.6))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = n
        set_run(r, size=14, bold=True, color=INK)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = lab
        set_run(r2, size=9, bold=True, color=GREEN)
    cta = pill(s, Inches(0.7), Inches(6.0), Inches(3.6), Inches(0.45), "Banco de Chile × Desafío Latam →", fill=INK, font_size=12)
    footer(s, 1)

    # 02 Agenda
    s = new_slide(prs)
    eyebrow(s, "AGENDA")
    title_block(s, ["Lo que verás"], "hoy", Inches(1.4))
    items = [
        ("03", "El desafío"),
        ("04", "La respuesta del proyecto"),
        ("05", "Impacto en cifras"),
        ("06", "Resultados académicos"),
        ("07", "Impacto en las personas"),
        ("08", "Inclusión y ecosistema"),
        ("09", "Aprendizajes clave"),
        ("10", "Proyección y cierre"),
    ]
    for i, (num, label) in enumerate(items):
        col = i % 2
        row = i // 2
        left = Inches(0.7) + Inches(col * 6.1)
        top = Inches(2.5) + Inches(row * 0.85)
        c = rect(s, left, top, Inches(5.7), Inches(0.7), fill=GRAY)
        t = s.shapes.add_textbox(left + Inches(0.25), top + Inches(0.18), Inches(5.2), Inches(0.4))
        tf = t.text_frame
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = num + "   "
        set_run(r1, size=14, bold=True, color=GREEN)
        r2 = p.add_run()
        r2.text = label
        set_run(r2, size=16, bold=True, color=INK)
    footer(s, 2)

    # 03 Desafío
    s = new_slide(prs)
    eyebrow(s, "EL DESAFÍO")
    y = title_block(s, ["Hay que preparar"], "el talento TP", Inches(2.8))
    lead = s.shapes.add_textbox(Inches(0.7), y, Inches(12), Inches(0.8))
    add_text(
        lead,
        [
            "La educación TP forma a jóvenes que se insertan temprano al trabajo, pero arrastra una brecha de competencias digitales frente a un mercado con IA, automatización y servicios financieros digitales."
        ],
        size=15,
        color=MUTED,
    )
    cards = [
        ("Brecha", "Estudiantes de 3° y 4° medio TP egresan sin dominio de herramientas digitales, IA aplicada ni comprensión del sistema financiero."),
        ("Riesgo", "Jóvenes en desventaja competitiva desde el primer día, y una modalidad percibida como “segunda categoría”."),
        ("Oportunidad", "Una intervención acotada, aplicada y certificada por Banco de Chile puede cerrar la brecha antes del egreso."),
    ]
    for i, (t, b) in enumerate(cards):
        card(s, Inches(0.7) + Inches(i * 4.1), Inches(3.5), Inches(3.9), Inches(2.4), t, [b])
    footer(s, 3)

    # 04 Respuesta
    s = new_slide(prs)
    eyebrow(s, "LA RESPUESTA")
    title_block(s, ["Itinerario doble"], "202 horas", Inches(2.5))
    lead = s.shapes.add_textbox(Inches(0.7), Inches(2.35), Inches(12), Inches(0.55))
    add_text(
        lead,
        ["Banco de Chile y Desafío Latam diseñaron un itinerario complementario: e-learning + presencial en el propio liceo."],
        size=15,
        color=MUTED,
    )
    card(
        s,
        Inches(0.7),
        Inches(3.1),
        Inches(5.9),
        Inches(2.9),
        "Transformación Digital con IA y Automatización",
        ["100% asincrónico y autoinstructivo.", "Franquicia Tributaria SENCE (precontrato)."],
        badge="E-LEARNING · 150 H",
    )
    card(
        s,
        Inches(6.9),
        Inches(3.1),
        Inches(5.9),
        Inches(2.9),
        "Cliente Digital en el Centro",
        ["6 módulos a medida.", "Bootcamp con ≥60% de práctica, casos reales y simulaciones."],
        badge="PRESENCIAL · 52 H",
    )
    footer(s, 4)

    # 05 Cifras
    s = new_slide(prs)
    eyebrow(s, "IMPACTO EN CIFRAS")
    title_block(s, ["Resultados del"], "primer ciclo", Inches(2.8))
    metrics = [
        ("85", "Estudiantes formados"),
        ("202", "Horas por estudiante"),
        ("92,5%", "Tasa de aprobación"),
        ("86,2%", "Asistencia promedio"),
    ]
    for i, (n, lab) in enumerate(metrics):
        metric(s, Inches(0.7) + Inches(i * 3.15), Inches(2.6), Inches(3.0), Inches(1.7), n, lab)
    card(
        s,
        Inches(0.7),
        Inches(4.6),
        Inches(5.9),
        Inches(1.9),
        "Cobertura",
        ["Colegio Initec · Curicó (Maule).", "Primer ciclo de un convenio proyectado a 1.660 estudiantes en 2026."],
    )
    card(
        s,
        Inches(6.9),
        Inches(4.6),
        Inches(5.9),
        Inches(1.9),
        "Satisfacción docente",
        ["10/10 global.", "4,4/5 equipo administrativo ADL · 4,2/5 experiencia como facilitador."],
    )
    footer(s, 5)

    # 06 Académico
    s = new_slide(prs)
    eyebrow(s, "DESEMPEÑO ACADÉMICO")
    title_block(s, ["Cuando se presentaron,"], "aprendieron", Inches(3.0))
    metric(s, Inches(0.7), Inches(2.7), Inches(3.5), Inches(1.6), "8,66", "Nota promedio presencial")
    metric(s, Inches(0.7), Inches(4.5), Inches(3.5), Inches(1.6), "9,66", "Quizzes e-learning")
    card(
        s,
        Inches(4.5),
        Inches(2.7),
        Inches(8.1),
        Inches(3.4),
        "Lectura clave",
        [
            "• Ningún estudiante reprobó por bajo rendimiento (6 casos por actividades pendientes).",
            "• Aprobación: 97,4% en 4° medio · 87,8% en 3° medio (base 80).",
            "• 60% del grupo con nota final entre 9,0 y 10,0.",
            "• Proyecto Integrador: 9,83 (3°) y 9,67 (4°).",
            "• E-learning: 70,4% finalización (69/98) · 73,8% avance promedio.",
        ],
    )
    footer(s, 6)

    # 07 Personas
    s = new_slide(prs)
    eyebrow(s, "IMPACTO EN LAS PERSONAS")
    title_block(s, ["Competencias para el"], "primer empleo", Inches(3.0))
    cards = [
        ("Digitales", "IA generativa (prompts CRTF, zero-shot / few-shot) y automatización no-code (Zapier, n8n, Make, Power Automate)."),
        ("Financieras", "Sistema financiero chileno y reguladores, análisis con hojas de cálculo, venta consultiva y asesoría a clientes."),
        ("Personales", "Presupuesto, clasificación de gastos, ahorro y planes de salida de deudas con cifras concretas."),
    ]
    for i, (t, b) in enumerate(cards):
        card(s, Inches(0.7) + Inches(i * 4.1), Inches(2.7), Inches(3.9), Inches(2.6), t, [b])
    quote_box(
        s,
        Inches(0.7),
        Inches(5.55),
        Inches(12),
        "Cada estudiante egresa con diploma Banco de Chile × Desafío Latam y un proyecto integrador ante cliente simulado.",
    )
    footer(s, 7)

    # 08 Inclusión
    s = new_slide(prs)
    eyebrow(s, "INCLUSIÓN · ECOSISTEMA")
    title_block(s, ["Articulación que"], "escala", Inches(1.8))
    card(
        s,
        Inches(0.7),
        Inches(2.5),
        Inches(5.9),
        Inches(3.5),
        "Inclusión y equidad",
        [
            "• Foco en liceos TP, segmento históricamente subvalorado.",
            "• Descentralización: Curicó, Región del Maule.",
            "• 100% vía SENCE: sin barrera económica.",
            "• Seguimiento caso a caso de rezagados.",
        ],
    )
    card(
        s,
        Inches(6.9),
        Inches(2.5),
        Inches(5.9),
        Inches(3.5),
        "Comunidades y validación",
        [
            "• Banca + OTEC digital + educación TP + Estado (SENCE).",
            "• Integrado al calendario escolar del liceo.",
            "• Certificación conjunta Banco de Chile – Desafío Latam.",
            "• Hitos con autoridades de las tres instituciones.",
        ],
    )
    footer(s, 8)

    # 09 Aprendizajes
    s = new_slide(prs)
    eyebrow(s, "APRENDIZAJES CLAVE")
    title_block(s, ["Qué funcionó · qué"], "ajustar", Inches(2.0))
    card(
        s,
        Inches(0.7),
        Inches(2.5),
        Inches(5.9),
        Inches(3.7),
        "Lo que funcionó",
        [
            "• Complementariedad e-learning + presencial.",
            "• Metodología práctica (≥60% de aula).",
            "• Coordinación sólida entre las 3 instituciones.",
            "• Seguimiento individualizado de rezagados.",
        ],
        badge="OK",
    )
    card(
        s,
        Inches(6.9),
        Inches(2.5),
        Inches(5.9),
        Inches(3.7),
        "Ajustar a futuro",
        [
            "• Anticipar calendario (asistencia 3° medio cayó a 17,9% al cierre).",
            "• Hitos intermedios en e-learning (29/98 no completaron).",
            "• Corregir rúbricas/formatos evaluativos.",
            "• Dinamizar Módulo 3 · mejorar LMS · encuesta a tiempo.",
        ],
        badge="NEXT",
    )
    footer(s, 9)

    # 10 Cierre
    s = new_slide(prs)
    eyebrow(s, "PROYECCIÓN · CIERRE")
    title_block(s, ["Listos para"], "escalar", Inches(2.0))
    lead = s.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(7), Inches(0.7))
    add_text(
        lead,
        ["Convenio 2026: 1.660 estudiantes — 1.500 en e-learning masivo y 160 en presencial (4 cohortes)."],
        size=16,
        color=MUTED,
    )
    body = s.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(7), Inches(2.0))
    add_text(
        body,
        [
            "• Replicar en Initec con calendario corregido.",
            "• Extender a nuevas especialidades y liceos.",
            "• Webinars BChile, empleabilidad y comunidad virtual.",
            "• Modelo SENCE validado operativamente.",
        ],
        size=15,
        color=MUTED,
    )
    quote_box(
        s,
        Inches(0.7),
        Inches(5.3),
        Inches(7),
        "Detrás de cada cifra hay un joven de Curicó que hoy usa IA, entiende el sistema financiero y llega a su primera entrevista con un diploma del Banco de Chile.",
    )
    panel = rect(s, Inches(8.2), Inches(2.3), Inches(4.4), Inches(4.0), fill=INK)
    pill(s, Inches(8.5), Inches(2.6), Inches(1.2), Inches(0.32), "NEXT", font_size=10)
    t = s.shapes.add_textbox(Inches(8.5), Inches(3.2), Inches(3.9), Inches(1.2))
    add_text(t, ["1.660", "estudiantes 2026"], size=28, bold=True, color=WHITE)
    # fix second line color - redo
    t2 = s.shapes.add_textbox(Inches(8.5), Inches(4.5), Inches(3.9), Inches(1.0))
    tf = t2.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "De un ciclo validado a la política de formación digital TP."
    set_run(r, size=14, color=GREEN_SOFT)
    pill(s, Inches(8.5), Inches(5.6), Inches(3.2), Inches(0.45), "Escalar el modelo →", fill=GREEN, font_size=13)
    footer(s, 10)

    prs.save(OUT)
    print(f"OK → {OUT}")


if __name__ == "__main__":
    build()
