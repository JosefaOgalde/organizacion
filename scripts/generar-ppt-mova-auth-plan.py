#!/usr/bin/env python3
"""Genera PPT: MOVA · Plan de ejecución — Login unificado mova_auth (12 pasos)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "index/clientes/mkof/MOVA-Auth-Plan-Ejecucion.pptx"

C_BG = RGBColor(0xE8, 0xF6, 0xF8)
C_ACCENT = RGBColor(0x4A, 0x7A, 0x80)
C_ACCENT_DARK = RGBColor(0x2A, 0x4A, 0x4E)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT = RGBColor(0x1F, 0x23, 0x28)
C_MUTED = RGBColor(0x65, 0x6D, 0x76)
C_HIGHLIGHT = RGBColor(0xFF, 0xF8, 0xC5)
C_OK = RGBColor(0x1A, 0x7F, 0x37)
C_WARN = RGBColor(0x9A, 0x67, 0x00)
C_BORDER = RGBColor(0xD0, 0xD7, 0xDE)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_header_bar(slide, title, subtitle=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()
    add_textbox(slide, Inches(0.55), Inches(0.18), Inches(11), Inches(0.5), title, size=24, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.62), Inches(11), Inches(0.35), subtitle, size=13, color=C_BG)


def add_step_badge(slide, num, left=Inches(0.55), top=Inches(1.35)):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = C_ACCENT
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p.alignment = PP_ALIGN.CENTER


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_ACCENT_DARK)
    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.6), "MOVA · Post-auditoría", size=20, color=RGBColor(0xA8, 0xD8, 0xDC))
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.0), "Plan de ejecución", size=44, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.3), Inches(11.5), Inches(0.9), "Login unificado · mova_auth", size=30, color=C_BG)
    add_textbox(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.5), "Hitos 2.1 + 2.2 · Semanas 1–2 del Gantt", size=16, color=RGBColor(0xA8, 0xD8, 0xDC))
    add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.4), "GRUPO MAKING OF · Guía para el encargado · Jul 2026", size=14, color=C_MUTED)


def slide_contexto(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "¿Qué vamos a lograr?", "Sin migrar de GoDaddy")
    bullets = [
        "Un solo login para todos los módulos M (MAESTRO, INGRESOS, EGRESOS…).",
        "mova_auth como único validador — ningún módulo valida por su cuenta.",
        "Sesión en cookie httpOnly + Secure (PHP en GoDaddy), no JWT en localStorage.",
        "Duración estimada: 8–10 días hábiles (Fases A → D).",
        "Prerrequisito: inventario hecho antes de tocar código en producción.",
    ]
    y = Inches(1.5)
    for b in bullets:
        add_textbox(slide, Inches(0.9), y, Inches(11.5), Inches(0.55), f"•  {b}", size=19, color=C_TEXT)
        y += Inches(0.62)
    warn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.0), Inches(11.8), Inches(1.5))
    warn.fill.solid()
    warn.fill.fore_color.rgb = C_HIGHLIGHT
    warn.line.color.rgb = RGBColor(0xD4, 0xA7, 0x2C)
    add_textbox(slide, Inches(0.95), Inches(5.2), Inches(11.3), Inches(1.1),
                "⚠  El problema no es falta de validación — es fragmentación:\nGoogle + mova_auth + permisos parciales + JWT en el navegador.", size=16, color=C_ACCENT_DARK)


def slide_mapa_fases(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Mapa de fases (12 pasos)", "Seguir en orden — un módulo a la vez")
    fases = [
        ("A", "Inventario", "Día 1", "Pasos 1–2 · Listar módulos y flujos actuales"),
        ("B", "Diseñar mova_auth", "Días 2–3", "Pasos 3–5 · Carpeta + archivos PHP núcleo"),
        ("C", "Login y sesión", "Días 4–7", "Pasos 6–9 · Cookie, validate.php, migrar módulos"),
        ("D", "Limpieza y cierre", "Días 8–10", "Pasos 10–12 · Quitar JWT, pruebas, documentar"),
    ]
    y = Inches(1.45)
    for letra, nombre, dias, desc in fases:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), y, Inches(11.8), Inches(1.15))
        box.fill.solid()
        box.fill.fore_color.rgb = C_WHITE
        box.line.color.rgb = C_ACCENT
        add_textbox(slide, Inches(0.95), y + Inches(0.12), Inches(1.2), Inches(0.4), f"Fase {letra}", size=14, bold=True, color=C_ACCENT)
        add_textbox(slide, Inches(2.1), y + Inches(0.08), Inches(5), Inches(0.4), nombre, size=18, bold=True, color=C_ACCENT_DARK)
        add_textbox(slide, Inches(7.5), y + Inches(0.1), Inches(2.5), Inches(0.35), dias, size=13, color=C_MUTED)
        add_textbox(slide, Inches(2.1), y + Inches(0.5), Inches(9.5), Inches(0.5), desc, size=14, color=C_TEXT)
        y += Inches(1.28)


def slide_paso(prs, num, fase, titulo, texto, destacar, tip, checklist):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, f"Paso {num} de 12 · Fase {fase}", titulo)
    add_step_badge(slide, num)
    add_textbox(slide, Inches(1.25), Inches(1.25), Inches(11), Inches(0.5), titulo, size=22, bold=True, color=C_ACCENT_DARK)
    add_textbox(slide, Inches(0.75), Inches(1.85), Inches(6.2), Inches(1.6), texto, size=16, color=C_TEXT)

    hl = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.85), Inches(5.5), Inches(1.0))
    hl.fill.solid()
    hl.fill.fore_color.rgb = C_HIGHLIGHT
    hl.line.color.rgb = RGBColor(0xD4, 0xA7, 0x2C)
    add_textbox(slide, Inches(7.4), Inches(2.0), Inches(5.1), Inches(0.75), f"👆 {destacar}", size=13, bold=True, color=C_ACCENT_DARK)

    tip_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.65), Inches(12.0), Inches(0.9))
    tip_box.fill.solid()
    tip_box.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFA)
    tip_box.line.color.rgb = C_BORDER
    add_textbox(slide, Inches(0.9), Inches(3.8), Inches(11.6), Inches(0.65), f"💡 Tip: {tip}", size=13, color=C_MUTED)

    add_textbox(slide, Inches(0.75), Inches(4.75), Inches(11.5), Inches(0.35), "Checklist de este paso:", size=14, bold=True, color=C_ACCENT)
    y = Inches(5.15)
    for item in checklist:
        add_textbox(slide, Inches(0.95), y, Inches(11), Inches(0.35), f"☐  {item}", size=13, color=C_TEXT)
        y += Inches(0.38)


PASOS = [
    ("A", "Listar todos los módulos M",
     "En cPanel → Administrador de archivos → public_html. Anota cada carpeta MOVA y su URL.",
     "Tabla: Módulo | URL | ¿Cómo valida hoy? | ¿localStorage?",
     "No tocar código aún — solo inventario.",
     ["Tabla de módulos completa", "URLs anotadas", "Responsable por módulo"]),
    ("A", "Documentar flujo actual",
     "Por módulo: ¿Google OAuth? ¿mova_auth? ¿JWT en localStorage? ¿llama n8n?",
     "Marcar en rojo lo que NO pase por mova_auth",
     "Dibujar o anotar el flujo actual antes de cambiar nada.",
     ["Flujo actual documentado", "JWT/localStorage identificados", "Endpoints n8n listados"]),
    ("B", "Definir mova_auth como único validador",
     "Acuerdo de equipo: ningún módulo valida permisos solo. Todos usan guard.php.",
     "Si no pasó por mova_auth → redirect a login",
     "Escribir la regla y compartirla con el equipo antes de subir archivos.",
     ["Regla escrita y compartida", "Lista de archivos acordada", "Responsable técnico asignado"]),
    ("B", "Crear carpeta mova_auth en GoDaddy",
     "cPanel → public_html → Nueva carpeta: mova_auth. Permisos 755.",
     "public_html/mova_auth/",
     "Probar URL https://tudominio.cl/mova_auth/ antes de seguir.",
     ["Carpeta creada", "Permisos 755", "URL responde en navegador"]),
    ("B", "Subir archivos núcleo PHP",
     "config.php, session.php, login.php, validate.php, guard.php, logout.php.",
     "guard.php se incluye AL INICIO de cada módulo",
     "config.php con secretos — no subir a repos públicos.",
     ["6 archivos PHP creados", "login.php abre en navegador", "guard.php probado en sandbox"]),
    ("C", "Implementar login único",
     "login.php: usuario entra una vez. Google opcional → resultado en sesión PHP, no localStorage.",
     "login.php → sesión → redirect al módulo pedido",
     "Reutilizar tokeninfo + whitelist como hoy en n8n.",
     ["Login funciona", "Redirect correcto", "Google tokeninfo probado"]),
    ("C", "Token de sesión propio (cookie)",
     "session_start con cookie httpOnly + Secure. El navegador envía la cookie sola.",
     "DevTools → Cookie HttpOnly ✓ — sin JWT en localStorage",
     "Requiere HTTPS (Cloudflare SSL Full).",
     ["Cookie HttpOnly activa", "Sin JWT en localStorage", "Sesión persiste al recargar"]),
    ("C", "Endpoint validate.php",
     "JSON { ok, usuario, permisos } si sesión válida. AJAX con credentials: include.",
     "Cada petición AJAX → validate.php primero",
     "401 sin sesión; 200 con datos de usuario.",
     ["validate.php 200 con sesión", "401 sin sesión", "Permisos por rol correctos"]),
    ("C", "Migrar módulos uno por uno",
     "Orden: sandbox primero, luego el resto. require guard.php al inicio de cada index.php.",
     "Un módulo a la vez — probar antes del siguiente",
     "Eliminar validación duplicada del módulo al migrar.",
     ["Sandbox migrado", "Módulos críticos migrados", "Todos los M usan guard.php"]),
    ("D", "Retirar JWT de localStorage",
     "Buscar localStorage con token/jwt en cada módulo. Eliminar — la cookie basta.",
     "DevTools → Application → Local Storage sin tokens",
     "Revisar también sessionStorage.",
     ["localStorage sin JWT", "Código viejo eliminado", "Sin regresiones"]),
    ("D", "Pruebas por módulo",
     "4 pruebas: sin login → redirect; con login → entra; logout; incógnito pide login.",
     "Chrome normal + ventana incógnito",
     "Repetir en cada módulo M migrado.",
     ["4 pruebas OK por módulo", "Logout funciona", "Sesión expira según config"]),
    ("D", "Documentar y cerrar hitos 2.1 + 2.2",
     "Actualizar ficha MOVA: URL login, config.php, acceso cPanel. Marcar Gantt.",
     "Entregable: mova_auth operativo + todos los M centralizados",
     "Capacitar al equipo en el nuevo flujo de login.",
     ["Documentación actualizada", "Equipo capacitado", "Hitos 2.1 y 2.2 cerrados en Gantt"]),
]


def slide_errores(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Errores frecuentes", "Si pasa esto, revisa esto")
    errores = [
        ("Redirect infinito", "login.php NO debe incluir guard.php"),
        ("Cookie no se guarda", "HTTPS activo — secure => true requiere SSL"),
        ("Módulo sigue pidiendo login", "Misma cookie path=/ y misma carpeta mova_auth"),
        ("AJAX devuelve 401", 'fetch con credentials: "include"'),
        ("Google OK pero no persiste", "session_start() antes de cualquier output"),
    ]
    y = Inches(1.5)
    for err, sol in errores:
        add_textbox(slide, Inches(0.9), y, Inches(4.5), Inches(0.4), err, size=15, bold=True, color=C_WARN)
        add_textbox(slide, Inches(5.5), y, Inches(7), Inches(0.4), f"→ {sol}", size=14, color=C_TEXT)
        y += Inches(0.55)


def slide_no_hacer(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Qué NO hacer", "Posición del cliente en el playbook")
    items = [
        "No migrar de GoDaddy ahora (movimiento lateral).",
        "No validar permisos dentro de cada módulo.",
        "No guardar tokens en localStorage.",
        "No migrar todos los módulos el mismo día sin probar.",
    ]
    y = Inches(1.55)
    for item in items:
        add_textbox(slide, Inches(0.9), y, Inches(11), Inches(0.5), f"✗  {item}", size=18, color=C_TEXT)
        y += Inches(0.6)


def slide_checklist_final(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Checklist final", "Antes de dar por cerrado el hito")
    items = [
        "mova_auth es el único punto de login",
        "Todos los módulos M incluyen guard.php",
        "Cookie httpOnly + Secure activa",
        "Sin JWT en localStorage",
        "validate.php operativo",
        "Pruebas OK en todos los módulos",
        "Documentación en ficha MOVA",
    ]
    y = Inches(1.5)
    for item in items:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y, Inches(0.35), Inches(0.35))
        box.fill.background()
        box.line.color.rgb = C_OK
        box.line.width = Pt(2)
        add_textbox(slide, Inches(1.35), y - Inches(0.05), Inches(10.5), Inches(0.45), item, size=17, color=C_TEXT)
        y += Inches(0.58)


def slide_siguiente(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_ACCENT_DARK)
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.7), "Después del login", size=32, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.4),
                "Gantt Semana 2–3:\n\n• Hito 3.1 — Modelar mova_datos (MySQL)\n• Hito 3.2 — MySQL fuente → Sheets vista\n• Hito 4.0 — Rutinas de operación",
                size=22, color=C_BG)
    add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.6),
                "Guía detallada: mova-auth-guia.html · PDF MOVA-Auth-Login-Unificado", size=16, color=RGBColor(0xA8, 0xD8, 0xDC))


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_contexto(prs)
    slide_mapa_fases(prs)
    for i, (fase, titulo, texto, dest, tip, chk) in enumerate(PASOS, start=1):
        slide_paso(prs, i, fase, titulo, texto, dest, tip, chk)
    slide_errores(prs)
    slide_no_hacer(prs)
    slide_checklist_final(prs)
    slide_siguiente(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Generado: {OUT}")
    print(f"Diapositivas: {len(prs.slides)}")


if __name__ == "__main__":
    main()
