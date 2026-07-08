#!/usr/bin/env python3
"""Genera PPT: MOVA · Día 1 — Status inventario módulos (acme-chile.cl)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "index/clientes/mkof/MOVA-D1-Inventario-Status.pptx"

C_BG = RGBColor(0xE8, 0xF6, 0xF8)
C_ACCENT = RGBColor(0x4A, 0x7A, 0x80)
C_ACCENT_DARK = RGBColor(0x2A, 0x4A, 0x4E)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT = RGBColor(0x1F, 0x23, 0x28)
C_MUTED = RGBColor(0x65, 0x6D, 0x76)
C_OK = RGBColor(0x1A, 0x7F, 0x37)
C_WARN = RGBColor(0xBF, 0x3F, 0x00)
C_HIGHLIGHT = RGBColor(0xFF, 0xF8, 0xC5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MODULOS = [
    ("Portal MOVA", "/mova/", "Google OAuth", "sí · axon_chats", "no", "Producción"),
    ("mova_auth", "/mova_auth/login.php", "PHP correo+clave", "pendiente", "parcial", "No gate del panel"),
    ("MOVA ERP", "/mova/erp/", "Contraseña local", "pendiente", "no", "Login aparte"),
    ("AXON", "/axon/", "Sin login visible", "pendiente", "no", "Chat directo"),
    ("RRHH", "/rrhh/", "403 Forbidden", "n/a", "no", "Sin index público"),
    ("Documentos", "/documentos/", "Público", "no", "n/a", "Playbooks"),
]

CHECKLIST = [
    ("Acceso cPanel GoDaddy", True),
    ("Árbol public_html/acme-chile.cl/", True),
    ("Auth documentada (módulos revisados)", True),
    ("localStorage en /mova/ (axon_chats)", True),
    ("n8n — delegado equipo n8n", True),
    ("Compartir status con equipo técnico", False),
    ("Marcar tarea mkof/01 completada", False),
]


def set_slide_bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_header_bar(slide, title, subtitle=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()
    add_textbox(slide, Inches(0.55), Inches(0.18), Inches(11), Inches(0.5), title, size=24, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.62), Inches(11), Inches(0.35), subtitle, size=13, color=C_BG)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_ACCENT_DARK)
    add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.5), "MOVA · Día 1 — Status", size=20, color=RGBColor(0xA8, 0xD8, 0xDC))
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.0), "Inventario módulos M", size=44, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.1), Inches(11.5), Inches(0.7), "acme-chile.cl · public_html/acme-chile.cl/", size=24, color=C_BG)
    add_textbox(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.5), "Tarea organizador: index.html?tarea=mkof/01", size=16, color=RGBColor(0xA8, 0xD8, 0xDC))
    add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.5), "Actualizado: 8 jul 2026 · GRUPO MAKING OF", size=14, color=C_MUTED)
    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.6), "Estado: EN PROGRESO — listo para cierre formal", size=18, bold=True, color=C_HIGHLIGHT)


def slide_resumen(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Resumen ejecutivo", "Día 1 · solo inventario — sin tocar código")
    items = [
        "Ruta real del sitio: public_html/acme-chile.cl/ (no en raíz de public_html).",
        "Problema confirmado: login FRAGMENTADO (Google + mova_auth PHP + contraseña ERP).",
        "Panel MOVA (/mova/) usa Google OAuth — NO pasa por mova_auth.",
        "localStorage en acme-chile.cl: clave axon_chats (widget AXON). Sin jwt/token visible.",
        "n8n: pendiente del equipo que lo administra (no bloquea este status).",
        "/pruebas/ queda fuera de alcance — era etapa 1, producción es /mova/.",
    ]
    y = Inches(1.45)
    for b in items:
        add_textbox(slide, Inches(0.85), y, Inches(11.5), Inches(0.55), f"•  {b}", size=17, color=C_TEXT)
        y += Inches(0.58)


def slide_checklist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Checklist Día 1", "Criterio de cierre")
    y = Inches(1.5)
    for label, ok in CHECKLIST:
        mark = "✓" if ok else "○"
        col = C_OK if ok else C_WARN
        add_textbox(slide, Inches(0.85), y, Inches(0.5), Inches(0.4), mark, size=20, bold=True, color=col)
        add_textbox(slide, Inches(1.35), y, Inches(10.5), Inches(0.4), label, size=18, color=C_TEXT)
        y += Inches(0.52)
    done = sum(1 for _, ok in CHECKLIST if ok)
    add_textbox(slide, Inches(0.85), Inches(6.2), Inches(11), Inches(0.45),
                f"Avance: {done}/{len(CHECKLIST)} ítems · Faltan: compartir + marcar tarea", size=16, bold=True, color=C_ACCENT_DARK)


def slide_tabla(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Módulos revisados en navegador", "https://acme-chile.cl/")
    headers = ["Módulo", "URL", "Auth", "localStorage", "mova_auth"]
    xs = [Inches(0.6), Inches(2.5), Inches(5.8), Inches(8.2), Inches(10.5)]
    for i, h in enumerate(headers):
        add_textbox(slide, xs[i], Inches(1.35), Inches(2.2), Inches(0.35), h, size=12, bold=True, color=C_ACCENT_DARK)
    y = Inches(1.75)
    for mod, url, auth, jwt, mova, nota in MODULOS:
        add_textbox(slide, xs[0], y, Inches(1.8), Inches(0.32), mod, size=11, bold=True, color=C_TEXT)
        add_textbox(slide, xs[1], y, Inches(3.2), Inches(0.32), url, size=10, color=C_MUTED)
        add_textbox(slide, xs[2], y, Inches(2.3), Inches(0.32), auth, size=10, color=C_TEXT)
        add_textbox(slide, xs[3], y, Inches(2.1), Inches(0.32), jwt, size=10, color=C_TEXT)
        add_textbox(slide, xs[4], y, Inches(1.5), Inches(0.32), mova, size=10, color=C_TEXT)
        y += Inches(0.38)
    add_textbox(slide, Inches(0.6), Inches(6.0), Inches(12), Inches(0.5),
                "Conclusión: ningún módulo M revisado usa mova_auth como único validador.", size=14, bold=True, color=C_WARN)


def slide_tabla_explicacion(prs):
    """Slide visual para audiencia no técnica — qué significa el inventario."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Inventario explicado", "Para quienes no trabajan con código")

    # Leyenda de columnas (tarjetas)
    cols = [
        ("Módulo", "Cada app o sección del sitio (portal, ERP, AXON…)."),
        ("Auth", "Cómo entra el usuario: Google, clave, o sin login."),
        ("localStorage", "Datos guardados en el navegador del usuario."),
        ("mova_auth", "¿Pasa por el login unificado? sí · no · parcial"),
    ]
    x0 = Inches(0.55)
    card_w = Inches(3.05)
    for i, (tit, desc) in enumerate(cols):
        left = x0 + i * Inches(3.15)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.35), card_w, Inches(1.55))
        box.fill.solid()
        box.fill.fore_color.rgb = C_WHITE
        box.line.color.rgb = C_ACCENT
        add_textbox(slide, left + Inches(0.15), Inches(1.5), card_w - Inches(0.3), Inches(0.35), tit, size=13, bold=True, color=C_ACCENT_DARK)
        add_textbox(slide, left + Inches(0.15), Inches(1.9), card_w - Inches(0.3), Inches(0.9), desc, size=11, color=C_TEXT)

    # Diagrama: login fragmentado vs objetivo
    add_textbox(slide, Inches(0.55), Inches(3.15), Inches(5.5), Inches(0.35), "Hoy: varias puertas de entrada", size=14, bold=True, color=C_WARN)
    puertas = [("Google", 0.55), ("PHP clave", 2.05), ("Contraseña ERP", 3.55), ("Sin login", 5.05)]
    for label, x in puertas:
        door = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.55), Inches(1.35), Inches(0.75))
        door.fill.solid()
        door.fill.fore_color.rgb = RGBColor(0xFF, 0xE8, 0xD6)
        door.line.color.rgb = C_WARN
        add_textbox(slide, Inches(x), Inches(3.72), Inches(1.35), Inches(0.4), label, size=10, bold=True, color=C_WARN, align=PP_ALIGN.CENTER)

    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.15), Inches(3.75), Inches(0.55), Inches(0.35))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = C_ACCENT
    arrow.line.fill.background()

    add_textbox(slide, Inches(6.85), Inches(3.15), Inches(5.8), Inches(0.35), "Objetivo MOVA: una sola puerta", size=14, bold=True, color=C_OK)
    gate = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(3.55), Inches(4.5), Inches(0.75))
    gate.fill.solid()
    gate.fill.fore_color.rgb = RGBColor(0xD4, 0xED, 0xDA)
    gate.line.color.rgb = C_OK
    add_textbox(slide, Inches(7.2), Inches(3.68), Inches(4.5), Inches(0.5), "mova_auth  →  todos los módulos M", size=12, bold=True, color=C_OK, align=PP_ALIGN.CENTER)

    # Hallazgo principal
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(4.65), Inches(12.2), Inches(1.15))
    box.fill.solid()
    box.fill.fore_color.rgb = C_HIGHLIGHT
    box.line.color.rgb = C_ACCENT
    add_textbox(slide, Inches(0.75), Inches(4.82), Inches(11.8), Inches(0.85),
                "Hallazgo Día 1: ningún módulo revisado usa solo mova_auth.\n"
                "Eso confirma que el login está repartido y hay que unificarlo en los días 2–7.",
                size=14, bold=True, color=C_ACCENT_DARK)

    # Chips de estado
    estados = [("sí", C_OK), ("no", C_WARN), ("parcial", C_WARN), ("pendiente", C_MUTED), ("n/a", C_MUTED)]
    add_textbox(slide, Inches(0.55), Inches(6.05), Inches(2.5), Inches(0.3), "Lectura rápida:", size=12, bold=True, color=C_ACCENT_DARK)
    x_chip = Inches(2.2)
    for label, col in estados:
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_chip, Inches(5.95), Inches(1.35), Inches(0.42))
        chip.fill.solid()
        chip.fill.fore_color.rgb = C_WHITE
        chip.line.color.rgb = col
        add_textbox(slide, x_chip, Inches(6.02), Inches(1.35), Inches(0.3), label, size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
        x_chip += Inches(1.55)


def add_round_box(slide, left, top, width, height, fill, line=None, line_w=1):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line:
        box.line.color.rgb = line
        box.line.width = Pt(line_w)
    else:
        box.line.fill.background()
    return box


def add_chip(slide, left, top, label, width=Inches(1.05), height=Inches(0.34), fill=C_WHITE, line=C_MUTED, text_color=C_TEXT, size=8, bold=False):
    add_round_box(slide, left, top, width, height, fill, line)
    add_textbox(slide, left, top + Inches(0.05), width, height - Inches(0.08), label, size=size, bold=bold, color=text_color, align=PP_ALIGN.CENTER)


def add_connector(slide, x1, y1, x2, y2, color=C_ACCENT):
    conn = slide.shapes.add_connector(1, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(1.5)
    return conn


def slide_cpanel(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Estructura cPanel", "public_html/acme-chile.cl/ — mapa visual del hosting")

    # Raíz del sitio
    root_w = Inches(5.2)
    root_left = Inches(4.05)
    add_round_box(slide, root_left, Inches(1.28), root_w, Inches(0.58), C_ACCENT_DARK, C_ACCENT)
    add_textbox(slide, root_left, Inches(1.4), root_w, Inches(0.35), "public_html / acme-chile.cl /", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Conectores desde raíz
    add_connector(slide, root_left + root_w / 2, Inches(1.86), Inches(3.4), Inches(2.18))
    add_connector(slide, root_left + root_w / 2, Inches(1.86), Inches(9.9), Inches(2.18))
    add_connector(slide, root_left + root_w / 2, Inches(1.86), root_left + root_w / 2, Inches(2.05))

    # Otros módulos en raíz (chips grises)
    otros = ["admin", "axon", "documentos", "rrhh", "crm", "skill", "operaciones", "multimedia"]
    x_otros = Inches(0.55)
    for nombre in otros:
        add_chip(slide, x_otros, Inches(2.05), nombre, width=Inches(1.15), fill=RGBColor(0xF0, 0xF4, 0xF5), line=C_MUTED, text_color=C_MUTED, size=8)
        x_otros += Inches(1.22)
    add_textbox(slide, Inches(0.55), Inches(2.42), Inches(12.2), Inches(0.25), "Otros módulos en la raíz del sitio", size=9, color=C_MUTED)

    # Carpeta mova/ (destacada)
    mova_left = Inches(0.55)
    mova_top = Inches(2.75)
    mova_w = Inches(5.95)
    mova_h = Inches(3.55)
    add_round_box(slide, mova_left, mova_top, mova_w, mova_h, RGBColor(0xE8, 0xF6, 0xF8), C_ACCENT, 2)
    add_round_box(slide, mova_left + Inches(0.12), mova_top + Inches(0.12), mova_w - Inches(0.24), Inches(0.48), C_ACCENT)
    add_textbox(slide, mova_left + Inches(0.2), mova_top + Inches(0.18), mova_w - Inches(0.4), Inches(0.35), "mova/  ·  Portal y apps MOVA", size=13, bold=True, color=C_WHITE)

    mova_mods = [
        "agencia", "brief", "cotizador", "cuentas", "doc", "erp", "estudios",
        "facturas", "forecast", "negocios", "oc", "operacion", "seo", "strack",
    ]
    x_m = mova_left + Inches(0.2)
    y_m = mova_top + Inches(0.78)
    chip_w = Inches(1.05)
    for i, mod in enumerate(mova_mods):
        highlight = mod == "erp"
        add_chip(
            slide, x_m, y_m, mod,
            width=chip_w,
            fill=RGBColor(0xFF, 0xF8, 0xC5) if highlight else C_WHITE,
            line=C_WARN if highlight else C_ACCENT,
            text_color=C_WARN if highlight else C_ACCENT_DARK,
            size=8,
            bold=highlight,
        )
        x_m += Inches(1.12)
        if (i + 1) % 5 == 0:
            x_m = mova_left + Inches(0.2)
            y_m += Inches(0.42)
    add_textbox(slide, mova_left + Inches(0.2), mova_top + Inches(2.95), mova_w - Inches(0.4), Inches(0.45),
                "El ERP no está en la raíz del sitio.\nEstá dentro de mova/erp/", size=10, bold=True, color=C_WARN)

    # Carpeta mova_auth/ (destacada)
    auth_left = Inches(6.85)
    auth_top = Inches(2.75)
    auth_w = Inches(5.95)
    auth_h = Inches(3.55)
    add_round_box(slide, auth_left, auth_top, auth_w, auth_h, RGBColor(0xF3, 0xEE, 0xF8), RGBColor(0x9A, 0x7A, 0xB8), 2)
    add_round_box(slide, auth_left + Inches(0.12), auth_top + Inches(0.12), auth_w - Inches(0.24), Inches(0.48), RGBColor(0x9A, 0x7A, 0xB8))
    add_textbox(slide, auth_left + Inches(0.2), auth_top + Inches(0.18), auth_w - Inches(0.4), Inches(0.35),
                "mova_auth/  ·  Login unificado (objetivo MOVA)", size=13, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    auth_files = ["login.php", "auth.php", "config.php", "google_login.php", "logout.php", "panel.php", "setup.sql"]
    x_a = auth_left + Inches(0.25)
    y_a = auth_top + Inches(0.85)
    for i, f in enumerate(auth_files):
        add_chip(
            slide, x_a, y_a, f,
            width=Inches(1.55),
            fill=C_WHITE,
            line=RGBColor(0x9A, 0x7A, 0xB8),
            text_color=RGBColor(0x5A, 0x40, 0x80),
            size=8,
        )
        x_a += Inches(1.62)
        if (i + 1) % 3 == 0:
            x_a = auth_left + Inches(0.25)
            y_a += Inches(0.42)
    add_textbox(slide, auth_left + Inches(0.25), auth_top + Inches(2.55), auth_w - Inches(0.5), Inches(0.8),
                "Aquí vivirá el único login\npara todos los módulos M.", size=11, bold=True, color=RGBColor(0x5A, 0x40, 0x80))

    # Leyenda inferior
    add_textbox(slide, Inches(0.55), Inches(6.45), Inches(12.2), Inches(0.35),
                "Lectura rápida:  teal = MOVA  ·  violeta = login unificado  ·  amarillo = ERP (ubicación clave)",
                size=11, bold=True, color=C_ACCENT_DARK)


def slide_localstorage(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "DevTools · Local Storage", "Sesión en https://acme-chile.cl/mova/")
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5), "Origen acme-chile.cl", size=18, bold=True, color=C_ACCENT_DARK)
    add_textbox(slide, Inches(0.9), Inches(2.1), Inches(11), Inches(0.45), "Clave: axon_chats  →  historial chats widget AXON", size=17, color=C_TEXT)
    add_textbox(slide, Inches(0.8), Inches(2.9), Inches(11.5), Inches(0.5), "Origen accounts.google.com (terceros)", size=18, bold=True, color=C_ACCENT_DARK)
    add_textbox(slide, Inches(0.9), Inches(3.5), Inches(11), Inches(0.45), "Sesión Google OAuth — fuera del top-level site", size=17, color=C_TEXT)
    add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(0.9),
                "No se detectó clave jwt / token / access_token en acme-chile.cl.\nSesión principal: Google OAuth + cookies.", size=16, color=C_TEXT)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.5), Inches(11.8), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = C_HIGHLIGHT
    box.line.color.rgb = C_ACCENT
    add_textbox(slide, Inches(0.95), Inches(5.65), Inches(11.4), Inches(0.7),
                "Objetivo Día 2–7: unificar en mova_auth con cookie HttpOnly — sin JWT en cliente.", size=15, bold=True, color=C_ACCENT_DARK)


def slide_siguiente(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Próximo paso — Día 2", "mkof/02 · Acuerdo + diseño mova_auth")
    pasos = [
        "1. Documento «Reglas-mova_auth» — regla: si no pasó por mova_auth, no entra.",
        "2. Lista de 6 archivos PHP: config, session, login, validate, guard, logout.",
        "3. Elegir módulo sandbox para migración (Día 5) — sugerencia: módulo simple bajo /mova/.",
        "4. Solicitar al equipo n8n listado de webhooks por módulo (complemento inventario).",
        "5. Marcar Día 1 completado en organizador tras compartir este PPT/PDF.",
    ]
    y = Inches(1.5)
    for p in pasos:
        add_textbox(slide, Inches(0.85), y, Inches(11.5), Inches(0.55), p, size=17, color=C_TEXT)
        y += Inches(0.62)
    add_textbox(slide, Inches(0.85), Inches(5.8), Inches(11.5), Inches(0.5),
                "Entregables: MOVA-D1-Inventario-Status.pptx · MOVA-D1-Inventario-Status.pdf", size=14, color=C_MUTED)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_title(prs)
    slide_resumen(prs)
    slide_checklist(prs)
    slide_tabla(prs)
    slide_tabla_explicacion(prs)
    slide_cpanel(prs)
    slide_localstorage(prs)
    slide_siguiente(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"PPT generado: {OUT}")
    print(f"Tamaño: {OUT.stat().st_size // 1024} KB")


if __name__ == "__main__":
    main()
