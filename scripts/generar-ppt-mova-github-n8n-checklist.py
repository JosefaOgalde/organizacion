#!/usr/bin/env python3
"""Genera PPT: MOVA · GitHub + n8n — checklist paso a paso (10 jul 2026)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "index/clientes/mkof/MOVA-GitHub-N8n-Checklist.pptx"

C_BG = RGBColor(0xE8, 0xF6, 0xF8)
C_ACCENT = RGBColor(0x4A, 0x7A, 0x80)
C_ACCENT_DARK = RGBColor(0x2A, 0x4A, 0x4E)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT = RGBColor(0x1F, 0x23, 0x28)
C_MUTED = RGBColor(0x65, 0x6D, 0x76)
C_OK = RGBColor(0x1A, 0x7F, 0x37)
C_WARN = RGBColor(0xBF, 0x3F, 0x00)
C_HIGHLIGHT = RGBColor(0xFF, 0xF8, 0xC5)
C_GREEN_BG = RGBColor(0xD4, 0xED, 0xDA)
C_ORANGE_BG = RGBColor(0xFF, 0xE8, 0xD6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, bold=False, color=C_TEXT, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_header_bar(slide, title, subtitle="", badge=""):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()
    add_textbox(slide, Inches(0.55), Inches(0.18), Inches(9.5), Inches(0.5), title, size=24, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.62), Inches(9.5), Inches(0.35), subtitle, size=13, color=C_BG)
    if badge:
        b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.8), Inches(0.28), Inches(2.0), Inches(0.5))
        b.fill.solid()
        b.fill.fore_color.rgb = C_ACCENT_DARK
        b.line.fill.background()
        add_textbox(slide, Inches(10.8), Inches(0.36), Inches(2.0), Inches(0.35), badge, size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


def add_round_box(slide, left, top, width, height, fill, line=None, lw=1):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line:
        box.line.color.rgb = line
        box.line.width = Pt(lw)
    else:
        box.line.fill.background()
    return box


def slide_checklist(slide, items, start_y=Inches(1.45)):
    y = start_y
    for label in items:
        add_textbox(slide, Inches(0.85), y, Inches(0.45), Inches(0.38), "○", size=18, bold=True, color=C_WARN)
        add_textbox(slide, Inches(1.35), y, Inches(11.3), Inches(0.38), label, size=15, color=C_TEXT)
        y += Inches(0.48)


def slide_pasos(slide, pasos, start_y=Inches(1.4)):
    y = start_y
    for i, (tit, desc) in enumerate(pasos, 1):
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.05), Inches(0.38), Inches(0.38))
        circ.fill.solid()
        circ.fill.fore_color.rgb = C_ACCENT
        circ.line.fill.background()
        add_textbox(slide, Inches(0.7), y + Inches(0.1), Inches(0.38), Inches(0.28), str(i), size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_round_box(slide, Inches(1.2), y, Inches(11.5), Inches(0.72), C_WHITE, C_ACCENT)
        add_textbox(slide, Inches(1.35), y + Inches(0.06), Inches(11.2), Inches(0.28), tit, size=12, bold=True, color=C_ACCENT_DARK)
        add_textbox(slide, Inches(1.35), y + Inches(0.34), Inches(11.2), Inches(0.32), desc, size=10, color=C_MUTED)
        y += Inches(0.82)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C_ACCENT_DARK)
    add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.5), "MOVA · Hito 1.1", size=20, color=RGBColor(0xA8, 0xD8, 0xDC))
    add_textbox(slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(1.0), "GitHub + solicitud n8n", size=40, bold=True, color=C_WHITE)
    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.6), "Checklist paso a paso · 10 jul 2026", size=22, color=C_BG)
    add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.5), "Paso 1 cuenta → Paso 2 repo → Paso 3 solicitud n8n (tabla + JSON + capturas)", size=16, color=RGBColor(0xA8, 0xD8, 0xDC))
    add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.55), "GRUPO MAKING OF · acme-chile.cl · auditoría MOVA", size=14, color=C_MUTED)


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Orden del día", "Tres pasos + checklist por cada uno")
    bloques = [
        ("1", "Crear cuenta GitHub", "Correo general del equipo · plan Free · verificar email"),
        ("2", "Repo privado", "mova-n8n-workflows · Private · vacío · copiar URL"),
        ("3", "Solicitud n8n", "Tabla + JSON + capturas por workflow · github-n8n.html"),
    ]
    y = Inches(1.5)
    for cod, tit, desc in bloques:
        add_round_box(slide, Inches(0.75), y, Inches(0.65), Inches(0.65), C_ACCENT_DARK)
        add_textbox(slide, Inches(0.75), y + Inches(0.15), Inches(0.65), Inches(0.35), cod, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_round_box(slide, Inches(1.55), y, Inches(10.9), Inches(0.65), C_WHITE, C_ACCENT)
        add_textbox(slide, Inches(1.7), y + Inches(0.08), Inches(10.6), Inches(0.28), tit, size=14, bold=True, color=C_ACCENT_DARK)
        add_textbox(slide, Inches(1.7), y + Inches(0.36), Inches(10.6), Inches(0.25), desc, size=11, color=C_MUTED)
        y += Inches(0.85)
    add_round_box(slide, Inches(0.75), Inches(4.2), Inches(11.8), Inches(1.1), C_HIGHLIGHT, C_ACCENT_DARK, 2)
    add_textbox(slide, Inches(0.95), Inches(4.4), Inches(11.4), Inches(0.75),
                "Tiempo estimado: ~30 min GitHub + 5 min enviar solicitud n8n.\n"
                "El trabajo mova_auth (D2–D5) sigue en paralelo — n8n no lo bloquea.",
                size=13, color=C_TEXT)


def slide_a_pasos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Bloque A — Crear cuenta GitHub", "Paso 1 · github.com/signup", badge="A")
    slide_pasos(slide, [
        ("Abrir registro", "Ir a github.com/signup (ventana incógnito si ya hay otra sesión)."),
        ("Correo general", "Email del proyecto — no personal. Varios deben leer el código de verificación."),
        ("Contraseña segura", "Mínimo 15 caracteres. Guardar en gestor del equipo (1Password, Bitwarden)."),
        ("Username", "Ej. mova-infra — solo letras, números y guiones. Difícil de cambiar después."),
        ("Create account", "Resolver captcha → botón verde Create account."),
        ("Verificar email", "Código de 8 dígitos del correo (revisar spam)."),
        ("Plan Free", "Skip preguntas opcionales → Continue for free."),
        ("Anotar datos", "Correo, usuario @ y ubicación de la contraseña en ficha MOVA."),
    ])


def slide_a_checklist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Checklist · Bloque A — Cuenta GitHub", "Marcar al completar cada ítem", badge="A ✓")
    slide_checklist(slide, [
        "Correo general del proyecto definido y accesible por el equipo",
        "Cuenta creada en github.com/signup",
        "Correo verificado con código de 8 dígitos",
        "Plan Free activo (repos privados incluidos)",
        "Usuario y contraseña guardados en gestor compartido del equipo",
        "2FA planificado para activar cuando el equipo lo acuerde",
        "Datos anotados en ficha MOVA (correo + usuario @)",
    ])
    add_textbox(slide, Inches(0.85), Inches(5.6), Inches(11.5), Inches(0.4),
                "Guía detallada: github-cuenta.html · MOVA-GitHub-Paso1-Crear-Cuenta.pdf",
                size=12, color=C_MUTED)


def slide_b_pasos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Bloque B — Repo privado", "Paso 2 · github.com/new", badge="B")
    slide_pasos(slide, [
        ("Iniciar sesión", "github.com/login con la cuenta del Bloque A."),
        ("New repository", "Botón + arriba derecha → New repository (o github.com/new)."),
        ("Owner correcto", "Debe ser la cuenta MOVA creada en el Paso 1."),
        ("Nombre exacto", "Repository name → mova-n8n-workflows (minúsculas, guiones)."),
        ("Descripción", "Opcional: Respaldo de workflows n8n — proyecto MOVA."),
        ("Private", "Visibility → Private (candado). Nunca público — lógica sensible."),
        ("Repo vacío", "NO marcar README, .gitignore ni license."),
        ("Create + URL", "Create repository → copiar URL HTTPS y anotar en ficha MOVA."),
        ("Colaboradores", "Opcional: Settings → Collaborators → Add people (rol Write)."),
    ])


def slide_b_checklist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Checklist · Bloque B — Repo privado", "Marcar al completar cada ítem", badge="B ✓")
    slide_checklist(slide, [
        "Sesión iniciada con la cuenta del Bloque A",
        "Repositorio mova-n8n-workflows creado",
        "Visibilidad Private confirmada (candado visible)",
        "Repositorio vacío — sin README inicial",
        "URL del repo copiada y anotada en ficha MOVA",
        "Colaboradores del equipo invitados (si aplica)",
        "Equipo técnico avisado: repo listo para recibir backup n8n",
    ])
    add_textbox(slide, Inches(0.85), Inches(5.6), Inches(11.5), Inches(0.4),
                "Guía detallada: github-repo.html · Paso 3 → github-n8n.html",
                size=12, color=C_MUTED)


def slide_c1_pasos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Bloque C · Pedido 1 — Inventario n8n", "Tabla operativa para auditoría MOVA", badge="C1")
    add_textbox(slide, Inches(0.75), Inches(1.25), Inches(12), Inches(0.4),
                "Pedir al equipo que administra n8n — complementa el inventario D1 (columna «¿n8n?»).",
                size=14, color=C_TEXT)
    cols = "Workflow | Activo | Webhook URL | Módulo | Auth | Método | Sensibles | Responsable | Notas"
    add_round_box(slide, Inches(0.75), Inches(1.75), Inches(11.8), Inches(0.55), C_WHITE, C_ACCENT)
    add_textbox(slide, Inches(0.9), Inches(1.88), Inches(11.5), Inches(0.35), cols, size=10, bold=True, color=C_ACCENT_DARK)
    pasos = [
        "Listar todos los workflows activos que sirven a acme-chile.cl o módulos M.",
        "Por cada uno: URL completa del webhook en producción.",
        "Indicar módulo consumidor: Portal /mova/, ERP, AXON, RRHH, etc.",
        "Documentar si requiere auth (token, API key, header, whitelist IP o ninguna).",
        "Enviar tabla por correo o carpeta compartida — sin contraseñas en texto plano.",
    ]
    y = Inches(2.5)
    for i, p in enumerate(pasos, 1):
        add_textbox(slide, Inches(0.85), y, Inches(11.5), Inches(0.42), f"{i}.  {p}", size=12, color=C_TEXT)
        y += Inches(0.48)
    add_round_box(slide, Inches(0.75), Inches(5.0), Inches(11.8), Inches(0.85), C_GREEN_BG, C_OK, 2)
    add_textbox(slide, Inches(0.95), Inches(5.15), Inches(11.4), Inches(0.6),
                "Este pedido NO se reemplaza con JSON — necesitamos la tabla para mapear módulo ↔ webhook ↔ auth.",
                size=12, bold=True, color=C_OK)


def slide_c1_checklist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Checklist · Pedido 1 — Inventario n8n", "Tabla de webhooks en producción", badge="C1 ✓")
    slide_checklist(slide, [
        "Contacto del equipo n8n identificado (nombre + correo)",
        "Solicitud enviada con columnas acordadas (ver slide anterior)",
        "Workflows de acme-chile.cl / módulos M listados",
        "URL de cada webhook en producción recibida",
        "Módulo consumidor indicado por cada webhook",
        "Tipo de autenticación documentado (sin secretos en claro)",
        "Responsable por workflow asignado",
        "Tabla integrada al Inventario-MOVA-modulos.md",
    ])


def slide_c2_pasos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Bloque C · Pedido 2 — Export JSON", "Backup para repo mova-n8n-workflows", badge="C2")
    slide_pasos(slide, [
        ("Export por workflow", "En n8n: workflow → menú ⋯ → Download / Export → archivo .json"),
        ("Solo activos", "Exportar workflows en producción (activos), no borradores obsoletos."),
        ("Entrega en zip", "Carpeta comprimida con todos los JSON o push directo al repo GitHub."),
        ("Sin secretos por correo", "Los JSON pueden referenciar credenciales por nombre — no enviar passwords."),
        ("Primera carga al repo", "Subir al repo privado cuando el equipo tenga acceso Write."),
    ], start_y=Inches(1.35))
    add_round_box(slide, Inches(0.75), Inches(5.55), Inches(11.8), Inches(0.75), C_ORANGE_BG, C_WARN, 2)
    add_textbox(slide, Inches(0.95), Inches(5.7), Inches(11.4), Inches(0.5),
                "El JSON guarda la estructura del workflow. Las credenciales reales siguen en n8n.",
                size=12, bold=True, color=C_WARN)


def slide_c2_checklist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Checklist · Pedido 2 — Export JSON", "Respaldo en GitHub", badge="C2 ✓")
    slide_checklist(slide, [
        "Solicitud de export JSON incluida en el mismo correo al equipo n8n",
        "Archivos .json recibidos por cada workflow activo",
        "Verificado que no vienen contraseñas en texto plano adjuntas",
        "JSON almacenados en carpeta segura o subidos a mova-n8n-workflows",
        "Fecha del primer backup anotada en ficha MOVA",
        "Responsable del backup periódico definido (semanal sugerido)",
    ])


def slide_c3_pasos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Paso 3 · Pedido 3 — Capturas n8n", "Entender el flujo visual por workflow", badge="3")
    slide_pasos(slide, [
        ("Lista workflows", "Captura de n8n con workflows activos (ON) visibles."),
        ("Canvas completo", "Por workflow: todos los nodos — trigger → lógica → respuesta."),
        ("Nodo Webhook", "Panel del webhook: URL producción, método HTTP, path."),
        ("Validación auth", "Nodos IF/Code donde validan token o usuario (clave D5)."),
        ("Executions", "Opcional: ejecución Success reciente — tapar datos sensibles."),
        ("Tapar secretos", "Pedir que oculten API keys y tokens antes de enviar."),
    ], start_y=Inches(1.3))
    add_round_box(slide, Inches(0.75), Inches(5.65), Inches(11.8), Inches(0.7), C_GREEN_BG, C_OK, 2)
    add_textbox(slide, Inches(0.95), Inches(5.8), Inches(11.4), Inches(0.45),
                "Las capturas complementan tabla y JSON — no las reemplazan. Guía: github-n8n.html",
                size=12, bold=True, color=C_OK)


def slide_c3_checklist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Checklist · Paso 3 — Capturas n8n", "Por cada workflow activo", badge="3 ✓")
    slide_checklist(slide, [
        "Correo enviado con los 3 pedidos (tabla + JSON + capturas)",
        "Lista de capturas requeridas incluida en el correo",
        "Captura lista workflows recibida",
        "Captura canvas por workflow activo",
        "Captura nodo Webhook por workflow (URL coincide con tabla)",
        "Capturas de validación auth recibidas o «sin auth» documentado",
        "Secretos tapados en todas las capturas",
        "Entregables organizados en carpeta por workflow",
    ])


def slide_c4_pasos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Paso 3 · Contexto complementario", "Además de tabla, JSON y capturas", badge="+")
    items = [
        ("URL instancia n8n", "Ej. n8n.empresa.cl — dónde administran."),
        ("Admin / contacto", "Quién tiene acceso admin para cambios futuros."),
        ("Webhooks sin auth", "¿Algún endpoint público sin validación? — riesgo a documentar."),
        ("Validación de usuario", "¿Algún workflow valida sesión/token de usuario? (clave para D5)."),
        ("Backup automático", "¿Existe sync a GitHub o hay que configurarlo después del repo?"),
    ]
    y = Inches(1.4)
    for tit, desc in items:
        add_round_box(slide, Inches(0.75), y, Inches(2.4), Inches(0.7), C_ACCENT)
        add_textbox(slide, Inches(0.75), y + Inches(0.22), Inches(2.4), Inches(0.35), tit, size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_round_box(slide, Inches(3.35), y, Inches(9.2), Inches(0.7), C_WHITE, C_MUTED)
        add_textbox(slide, Inches(3.5), y + Inches(0.2), Inches(8.9), Inches(0.35), desc, size=12, color=C_TEXT)
        y += Inches(0.82)


def slide_c4_checklist(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Checklist · Contexto n8n", "Complemento al Paso 3", badge="+ ✓")
    slide_checklist(slide, [
        "URL de la instancia n8n recibida",
        "Contacto admin n8n confirmado",
        "Listado de webhooks sin autenticación (si existen)",
        "Workflows que validan usuario/sesión identificados",
        "Decisión sobre backup automático vs manual documentada",
    ])


def slide_no_pedir(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Qué NO pedir (o pedir con cuidado)", "Evitar riesgos de seguridad")
    rows = [
        ("Contraseñas / API keys por correo", "Pedir solo tipo de auth y nombre de credencial en n8n"),
        ("Solo JSON o solo capturas", "Pedir tabla + JSON + capturas (los tres)"),
        ("Acceso admin n8n hoy", "Solo listado + exports; admin puede ser después"),
        ("Export con secretos embebidos", "Export estándar; credenciales documentadas aparte en n8n"),
    ]
    y = Inches(1.45)
    for mal, bien in rows:
        add_round_box(slide, Inches(0.75), y, Inches(5.3), Inches(0.65), C_ORANGE_BG, C_WARN)
        add_textbox(slide, Inches(0.9), y + Inches(0.18), Inches(5.0), Inches(0.35), f"✗  {mal}", size=11, bold=True, color=C_WARN)
        add_round_box(slide, Inches(6.3), y, Inches(6.25), Inches(0.65), C_GREEN_BG, C_OK)
        add_textbox(slide, Inches(6.45), y + Inches(0.12), Inches(5.95), Inches(0.45), f"✓  {bien}", size=10, color=C_OK)
        y += Inches(0.78)


def slide_correo(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Texto listo para enviar al equipo n8n", "Copiar y pegar en correo")
    texto = (
        "Asunto: Solicitud MOVA — inventario y respaldo n8n (acme-chile.cl)\n\n"
        "Hola,\n\n"
        "Para la auditoría MOVA en acme-chile.cl necesitamos:\n\n"
        "1) TABLA de workflows en producción que sirvan al sitio o módulos MOVA:\n"
        "   · Nombre · URL webhook · Módulo · Auth · Responsable · Notas\n\n"
        "2) EXPORT JSON de esos workflows activos (backup repo privado mova-n8n-workflows).\n\n"
        "3) CAPTURAS por workflow activo:\n"
        "   · Lista workflows · canvas completo · nodo Webhook\n"
        "   · nodos validación auth · ejecución Success (opcional)\n"
        "   Tapar secretos antes de enviar.\n\n"
        "4) CONTEXTO: URL instancia n8n · contacto admin · webhooks sin auth.\n\n"
        "No enviar contraseñas ni tokens por correo.\n\n"
        "Plazo sugerido: [FECHA]\n"
        "Contacto MOVA: [TU CORREO]\n\n"
        "Gracias."
    )
    add_round_box(slide, Inches(0.75), Inches(1.35), Inches(11.8), Inches(5.5), C_WHITE, C_ACCENT, 2)
    add_textbox(slide, Inches(0.95), Inches(1.5), Inches(11.4), Inches(5.2), texto, size=11, color=C_TEXT)


def slide_cierre(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_header_bar(slide, "Checklist maestro — cierre del día", "GitHub + solicitud n8n enviada")
    slide_checklist(slide, [
        "Paso 1 — Cuenta GitHub creada y verificada",
        "Paso 2 — Repo mova-n8n-workflows privado · URL anotada",
        "Paso 3 — Correo n8n enviado (tabla + JSON + capturas)",
        "Capturas requeridas listadas en el correo",
        "Ficha MOVA actualizada (GitHub + repo + contacto n8n)",
        "Seguimiento agendado cuando responda el equipo n8n",
        "Inventario-MOVA-modulos.md listo para columna n8n",
        "Continuar mova_auth (D3 archivos núcleo) en paralelo",
    ], start_y=Inches(1.4))
    add_textbox(slide, Inches(0.85), Inches(5.85), Inches(11.5), Inches(0.45),
                "Guías: github-cuenta.html · github-repo.html · github-n8n.html",
                size=13, bold=True, color=C_ACCENT_DARK)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide_title(prs)
    slide_agenda(prs)
    slide_a_pasos(prs)
    slide_a_checklist(prs)
    slide_b_pasos(prs)
    slide_b_checklist(prs)
    slide_c1_pasos(prs)
    slide_c1_checklist(prs)
    slide_c2_pasos(prs)
    slide_c2_checklist(prs)
    slide_c3_pasos(prs)
    slide_c3_checklist(prs)
    slide_c4_pasos(prs)
    slide_c4_checklist(prs)
    slide_no_pedir(prs)
    slide_correo(prs)
    slide_cierre(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"PPT generado: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Tamaño: {OUT.stat().st_size // 1024} KB")


if __name__ == "__main__":
    main()
